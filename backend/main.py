"""FastAPI backend for LLM Council."""

import logging
import time
from fastapi import Depends, FastAPI, Header, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Dict, Any, Optional
import uuid
import json
import io
import os
import asyncio
import base64
import random
import re
from datetime import datetime
from starlette.middleware.base import BaseHTTPMiddleware

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

from . import storage
from .council import run_full_council, generate_conversation_title, build_conversation_context, stage1_collect_responses, stage2_collect_rankings, stage3_synthesize_final, build_stage3_prompt, calculate_aggregate_rankings, stage2_ca_validation_pass, doubting_thomas_review, parse_ranking_from_text, parse_rubric_scores, parse_claim_counts, compute_relevancy_gate
from .config import OPENROUTER_API_KEY, AVAILABLE_MODELS, DEFAULT_COUNCIL_MODELS, DEFAULT_CHAIRMAN_MODEL, COUNCIL_MODELS, GOOGLE_API_KEY, is_google_model
from .model_sync import sync_models, get_live_models, get_defaults, get_sync_status, periodic_sync_loop
from .openrouter import query_model, query_model_stream
from .resilience import (
    kill_switch,
    circuit_breaker,
    health_monitor,
    resolve_fallback,
    check_quorum,
    KillSwitchError,
    QuorumError,
    MIN_STAGE1_QUORUM,
    MIN_STAGE2_QUORUM,
)
from .grounding import compute_response_grounding_scores, get_rubric_criteria, enhance_ca_with_validation
from .skills import run_evidence_skills, format_citations_for_prompt
from .token_tracking import SessionCostTracker
from .pipeline_timer import PipelineTimer
from .memory import get_memory_manager, get_user_profile_memory, get_eca
from .memory_store import set_memory_user
from .prompt_guard import evaluate_prompt
from . import citation as citation_registry
from .orchestrator import (
    pre_stage1_agent,
    post_stage2_agent,
    post_stage3_agent,
    user_gate_agent,
)
from .agents import run_agent_team, enrich_stage3_citations, validate_and_fix_citations
from .openrouter import close_shared_client
from .security import get_security_status
from .infographics import extract_infographic, strip_infographic_block
from .auth import get_authenticated_user_id
from .config import ENTRA_SSO_ENABLED


from .health_probe import health_agent, periodic_health_check


def check_token_expiry():
    """Check JWT token expiration on startup."""
    if not OPENROUTER_API_KEY:
        print("⚠️  WARNING: No API key configured in .env!")
        return
    
    # Check if it's a persistent API key (mga-*) instead of JWT
    if OPENROUTER_API_KEY.startswith("mga-"):
        print(f"\n{'='*60}")
        print(f"🔐 API Key Status")
        print(f"{'='*60}")
        print(f"   Type: Persistent myGenAssist API Key (mga-*)")
        print(f"   Key: {OPENROUTER_API_KEY[:20]}...")
        print(f"   🟢 STATUS: OK (Persistent keys do not expire)")
        print(f"{'='*60}\n")
        return
    
    try:
        payload = OPENROUTER_API_KEY.split('.')[1]
        payload += '=' * (4 - len(payload) % 4)
        decoded = json.loads(base64.urlsafe_b64decode(payload))
        
        exp_date = datetime.fromtimestamp(decoded.get('exp', 0))
        remaining_mins = (exp_date - datetime.now()).total_seconds() / 60
        
        print(f"\n{'='*60}")
        print(f"🔐 JWT Token Status")
        print(f"{'='*60}")
        print(f"   Expires: {exp_date.strftime('%d-%b-%Y %H:%M:%S')}")
        print(f"   Remaining: {remaining_mins:.0f} minutes")
        
        if remaining_mins <= 0:
            print(f"   ⛔ STATUS: EXPIRED!")
            print(f"\n   Run: python token_monitor.py --refresh <new_token>")
        elif remaining_mins <= 10:
            print(f"   🔴 STATUS: CRITICAL - Token expiring soon!")
        elif remaining_mins <= 30:
            print(f"   🟠 STATUS: WARNING - Consider refreshing token")
        else:
            print(f"   🟢 STATUS: OK")
        print(f"{'='*60}\n")
        
    except Exception as e:
        print(f"⚠️  Could not decode JWT token: {e}")


# Check token on startup
check_token_expiry()


# ── Lifespan: startup model sync + periodic refresh ─────────────────────
from contextlib import asynccontextmanager

@asynccontextmanager
async def lifespan(application):
    """Run model sync on startup, launch periodic refresh in background."""
    logger.info("🚀 Running initial model sync...")
    try:
        summary = await asyncio.wait_for(sync_models(), timeout=25)
        logger.info(f"   Initial sync result: {summary.get('total_after_filter', '?')} models")
    except asyncio.TimeoutError:
        logger.warning("⚠️  Model sync timed out after 25s — starting with default models")
    except Exception as e:
        logger.warning(f"⚠️  Model sync failed: {e} — starting with default models")
    # Launch periodic sync as a background task
    sync_task = asyncio.create_task(periodic_sync_loop())
    # Launch health probe agent (checks every 5 minutes)
    health_task = asyncio.create_task(periodic_health_check(interval_seconds=300))
    yield
    health_task.cancel()
    sync_task.cancel()
    await close_shared_client()


app = FastAPI(title="LLM Council API", lifespan=lifespan)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Request Logging Middleware
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class RequestLoggingMiddleware(BaseHTTPMiddleware):
    """Log incoming requests with all headers except authorization (excluding /health endpoint)."""
    
    async def dispatch(self, request: Request, call_next):
        # Skip logging for health endpoint (used by ECS)
        if request.url.path == "/health":
            return await call_next(request)
        
        # Log endpoint
        logger.info(f"[Request] {request.method} {request.url.path}")
        
        # Log all headers except authorization
        filtered_headers = {}
        for header_name, header_value in request.headers.items():
            if header_name.lower() not in ["authorization", "auth", "token"]:
                filtered_headers[header_name] = header_value
        
        logger.info(f"[Request] Headers: {json.dumps(filtered_headers, indent=2)}")
        
        # Continue processing the request
        response = await call_next(request)
        return response


# Add request logging middleware
app.add_middleware(RequestLoggingMiddleware)

# Enable CORS — with credentials support for Entra ID SSO
_cors_origins = [
    "http://localhost:5173",
    "https://llmcouncil-frontend.azurewebsites.net",
    "https://llmcouncil-agents.ai",
]
app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors_origins if ENTRA_SSO_ENABLED else ["*"],
    allow_credentials=ENTRA_SSO_ENABLED,
    allow_methods=["*"],
    allow_headers=["*"],
)


class CreateConversationRequest(BaseModel):
    """Request to create a new conversation."""
    pass


class AttachmentData(BaseModel):
    """Attachment file data — supports both base64 inline and Azure Blob upload."""
    name: str
    type: str
    size: int
    base64: str = ""        # Inline content (legacy / local dev)
    blob_name: str = ""     # Azure Blob reference (cloud — SAS upload)


class SendMessageRequest(BaseModel):
    """Request to send a message in a conversation."""
    content: str
    attachments: List[AttachmentData] = []
    council_models: Optional[List[str]] = None
    chairman_model: Optional[str] = None
    web_search_enabled: bool = False
    speed_mode: bool = False


class EnhancePromptRequest(BaseModel):
    """Request to enhance a user prompt."""
    content: str


# Allowed MIME types for file attachments
ALLOWED_MIME_TYPES = {
    'application/pdf': 'PDF',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'PowerPoint',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'Excel',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'Word',
    'text/markdown': 'Markdown',
    'text/plain': 'Text',
    'image/png': 'Image (PNG)',
    'image/jpeg': 'Image (JPEG)',
    'image/gif': 'Image (GIF)',
    'image/webp': 'Image (WebP)',
    'image/svg+xml': 'Image (SVG)',
}

MAX_FILE_SIZE_BLOB = 200 * 1024 * 1024   # 200 MB via Azure Blob SAS upload
MAX_FILE_SIZE_INLINE = 10 * 1024 * 1024  # 10 MB via base64 JSON body


def validate_attachment(attachment: AttachmentData) -> str | None:
    """Validate an attachment. Returns error message if invalid, None if valid."""
    if attachment.type not in ALLOWED_MIME_TYPES:
        return f"Invalid file type: {attachment.type}. Allowed: PDF, PPTX, XLSX, DOCX"

    # Blob-uploaded files have a higher limit
    limit = MAX_FILE_SIZE_BLOB if attachment.blob_name else MAX_FILE_SIZE_INLINE
    if attachment.size > limit:
        limit_mb = limit // (1024 * 1024)
        return f"File too large: {attachment.name}. Maximum: {limit_mb}MB"

    return None


def extract_file_content_description(attachment: AttachmentData) -> str:
    """Extract actual text content from an attachment for LLM context.

    Supports two modes:
    - **Blob mode** (``blob_name`` set): downloads bytes from Azure Blob Storage
    - **Inline mode** (``base64`` set): decodes the base64 payload

    Falls back to a placeholder description when extraction fails.
    """
    import base64
    import io

    file_type = ALLOWED_MIME_TYPES.get(attachment.type, 'Document')
    fallback = f"[Attached {file_type} file: {attachment.name} ({attachment.size / 1024:.1f} KB) — content could not be extracted]"
    MAX_CHARS = 80_000  # ~20k tokens — keep context window manageable

    try:
        if attachment.blob_name:
            # ── Blob mode: download from Azure Blob Storage ──────
            from .storage import download_attachment_blob
            file_bytes = download_attachment_blob(attachment.blob_name)
            buf = io.BytesIO(file_bytes)
        else:
            # ── Inline mode: decode base64 payload ───────────────
            raw_b64 = attachment.base64
            if not raw_b64:
                return fallback
            if "," in raw_b64[:80]:
                raw_b64 = raw_b64.split(",", 1)[1]
            file_bytes = base64.b64decode(raw_b64)
            buf = io.BytesIO(file_bytes)
    except Exception as exc:
        logger.warning(f"[Attachment] data load failed for {attachment.name}: {exc}")
        return fallback

    extracted_text: str | None = None

    try:
        # ── PDF ──────────────────────────────────────────────────
        if attachment.type == "application/pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(buf)
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    pages.append(f"--- Page {i + 1} ---\n{text}")
            extracted_text = "\n\n".join(pages) if pages else None

        # ── DOCX ─────────────────────────────────────────────────
        elif attachment.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            from docx import Document
            doc = Document(buf)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            extracted_text = "\n\n".join(paragraphs) if paragraphs else None

        # ── PPTX ─────────────────────────────────────────────────
        elif attachment.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            from pptx import Presentation
            prs = Presentation(buf)
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
            extracted_text = "\n\n".join(slides) if slides else None

        # ── XLSX ─────────────────────────────────────────────────
        elif attachment.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            from openpyxl import load_workbook
            wb = load_workbook(buf, read_only=True, data_only=True)
            sheets = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    sheets.append(f"--- Sheet: {ws.title} ---\n" + "\n".join(rows))
            wb.close()
            extracted_text = "\n\n".join(sheets) if sheets else None

        # ── Plain Text / Markdown ────────────────────────────────
        elif attachment.type in ("text/plain", "text/markdown"):
            extracted_text = file_bytes.decode("utf-8", errors="replace")

    except Exception as exc:
        logger.warning(f"[Attachment] Text extraction failed for {attachment.name}: {exc}")
        return fallback

    if not extracted_text or not extracted_text.strip():
        return fallback

    # Truncate to keep context window manageable
    if len(extracted_text) > MAX_CHARS:
        extracted_text = extracted_text[:MAX_CHARS] + "\n\n[… content truncated at 80 000 characters]"

    return (
        f"=== Content extracted from {file_type} file: {attachment.name} "
        f"({attachment.size / 1024:.1f} KB) ===\n\n"
        f"{extracted_text}\n\n"
        f"=== End of {attachment.name} ==="
    )


async def get_user_id(user_id: str = Header(..., alias="user-id")) -> str:
    """Extract and validate the user-id header injected by the reverse proxy.
    
    DEPRECATED: Use get_authenticated_user_id from auth.py instead.
    Kept for backward compatibility with non-conversation endpoints.
    """
    sanitized = user_id.strip()
    if not sanitized or "/" in sanitized or "\\" in sanitized or ".." in sanitized:
        raise HTTPException(status_code=400, detail="Invalid user-id header")
    return sanitized


class ConversationMetadata(BaseModel):
    """Conversation metadata for list view."""
    id: str
    created_at: str
    title: str
    message_count: int
    context_tags: Optional[Dict[str, Any]] = None


class Conversation(BaseModel):
    """Full conversation with all messages."""
    id: str
    created_at: str
    title: str
    messages: List[Dict[str, Any]]
    context_tags: Optional[Dict[str, Any]] = None


@app.get("/")
async def root():
    """Health check endpoint."""
    return {"status": "ok", "service": "LLM Council API"}


@app.get("/health")
async def health():
    """Lightweight infra health check endpoint for ALB/ECS."""
    return {"status": "ok"}


@app.get("/api/health/deep")
async def deep_health():
    """Comprehensive health check — tests all subsystems (DB, API, memory, models, resilience)."""
    return await health_agent.run_deep_check()


@app.get("/api/health/history")
async def health_history(limit: int = 20):
    """Get recent health check history."""
    return {"history": health_agent.get_history(limit=limit)}


@app.get("/api/health/failures")
async def health_failures():
    """Get subsystems with consecutive failures."""
    return health_agent.get_failure_report()


# ── Azure Speech Token (for frontend Speech SDK) ────────────────────

_AZURE_SPEECH_KEY    = os.getenv("AZURE_SPEECH_KEY", "")
_AZURE_SPEECH_REGION = os.getenv("AZURE_SPEECH_REGION", "eastus")

@app.get("/api/speech/token")
async def speech_token():
    """
    Issue a short-lived authorization token for the Azure Speech SDK.

    The frontend calls this instead of embedding the Speech key directly.
    Tokens are valid for 10 minutes; the SDK auto-refreshes via this
    endpoint when needed.
    """
    if not _AZURE_SPEECH_KEY:
        raise HTTPException(
            status_code=503,
            detail="Azure Speech not configured — AZURE_SPEECH_KEY missing",
        )
    import httpx
    token_url = (
        f"https://{_AZURE_SPEECH_REGION}.api.cognitive.microsoft.com"
        "/sts/v1.0/issueToken"
    )
    try:
        async with httpx.AsyncClient(http2=True, verify=False, timeout=10) as client:
            resp = await client.post(
                token_url,
                headers={"Ocp-Apim-Subscription-Key": _AZURE_SPEECH_KEY},
            )
            resp.raise_for_status()
            return {"token": resp.text, "region": _AZURE_SPEECH_REGION}
    except Exception as e:
        logger.error(f"Speech token fetch failed: {e}")
        raise HTTPException(status_code=502, detail="Failed to fetch speech token")


# ── Citation Registry ────────────────────────────────────────────────

@app.get("/api/citations")
async def get_citations(module: Optional[str] = None, q: Optional[str] = None):
    """
    Return the citation registry for auditability / traceability.

    Query params:
        module — filter by backend module (e.g. 'council', 'grounding')
        q      — full-text search across title, abstract, relevance
    """
    if q:
        return {"citations": citation_registry.search(q)}
    if module:
        return {"citations": citation_registry.get_by_module(module)}
    return {"citations": citation_registry.list_all(), "stats": citation_registry.stats()}


@app.get("/api/citations/{citation_id}")
async def get_citation_detail(citation_id: str):
    """Return a single citation with APA and BibTeX formatted strings."""
    cite = citation_registry.get_citation(citation_id)
    if not cite:
        raise HTTPException(status_code=404, detail=f"Citation '{citation_id}' not found")
    return {
        "citation": cite,
        "formatted": {
            "apa": citation_registry.format_apa(citation_id),
            "bibtex": citation_registry.format_bibtex(citation_id),
        },
    }


@app.get("/api/models")
async def get_available_models():
    """Get list of available models — live from MyGenAssist API + Google AI Studio."""
    live = get_live_models()
    defaults = get_defaults()
    # Fallback to static config if sync hasn't populated yet
    if not live:
        from .config import get_all_available_models
        live = get_all_available_models()
        defaults = {
            "council_models": DEFAULT_COUNCIL_MODELS,
            "chairman_model": DEFAULT_CHAIRMAN_MODEL,
        }

    # Tag Bayer models with provider if missing
    all_models = [
        {**m, "provider": m.get("provider", "bayer")} for m in live
    ]

    return {
        "models": all_models,
        "defaults": defaults,
        "google_enabled": bool(GOOGLE_API_KEY),
    }


@app.get("/api/models/google")
async def discover_google_models():
    """Live discovery of available Google AI Studio models."""
    if not GOOGLE_API_KEY:
        return {"models": [], "error": "GOOGLE_API_KEY not configured in .env"}
    from .google_provider import list_google_models
    models = await list_google_models()
    return {"models": models, "count": len(models)}


@app.post("/api/models/sync")
async def trigger_model_sync():
    """Manually trigger a model sync from the MyGenAssist catalog."""
    summary = await sync_models()
    return summary


@app.get("/api/models/sync-status")
async def model_sync_status():
    """Get the current model sync status and last sync time."""
    return get_sync_status()


# ── Attachment Blob Upload (SAS token) ────────────────────────────────

class UploadUrlRequest(BaseModel):
    """Request a SAS URL for direct-to-blob attachment upload."""
    filename: str
    content_type: str
    size: int


class UploadUrlResponse(BaseModel):
    """SAS URL + blob reference returned to the frontend."""
    upload_url: str
    blob_name: str


@app.post("/api/attachments/upload-url", response_model=UploadUrlResponse)
async def get_upload_url(
    request: UploadUrlRequest,
    user_id: str = Depends(get_authenticated_user_id),
):
    """Generate a time-limited SAS URL for the browser to PUT a file directly
    into Azure Blob Storage, bypassing App Service body-size limits.
    """
    from .storage import is_blob_configured, generate_attachment_upload_url

    if not is_blob_configured():
        raise HTTPException(
            status_code=503,
            detail="Blob storage not configured — use inline base64 upload",
        )

    if request.content_type not in ALLOWED_MIME_TYPES:
        raise HTTPException(status_code=400, detail=f"Invalid file type: {request.content_type}")

    if request.size > MAX_FILE_SIZE_BLOB:
        raise HTTPException(
            status_code=400,
            detail=f"File too large. Maximum: {MAX_FILE_SIZE_BLOB // (1024*1024)}MB",
        )

    try:
        upload_url, blob_name = generate_attachment_upload_url(
            user_id=user_id,
            filename=request.filename,
            content_type=request.content_type,
        )
        logger.info(f"[SAS Upload] Generated URL for {request.filename} ({request.size} bytes) → {blob_name}")
        return UploadUrlResponse(upload_url=upload_url, blob_name=blob_name)
    except Exception as exc:
        logger.error(f"[SAS Upload] Failed to generate URL for {request.filename}: {exc}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"SAS URL generation failed: {exc}")


@app.post("/api/enhance-prompt")
async def enhance_prompt(request: EnhancePromptRequest):
    """
    Enhance a user's prompt to be more specific, detailed, and effective.
    Uses a fast model (gemini-2.5-flash) to generate an improved version.

    Follow-up aware: If the prompt targets a specific Stage or model
    (e.g. "Regarding Stage 2: ..."), the routing prefix is preserved
    verbatim and only the question body is enhanced.
    """
    if not request.content.strip():
        raise HTTPException(status_code=400, detail="Prompt content is required")

    # ── Detect and extract follow-up routing prefix ──────────────
    _RE_ENHANCE_STAGE = re.compile(
        rf"^({_FOLLOWUP_PREFIX}\s+Stage\s*\d){_SEPARATOR}(.+)",
        re.DOTALL | re.IGNORECASE,
    )
    _RE_ENHANCE_MODEL = re.compile(
        rf"^({_FOLLOWUP_PREFIX}\s+.+?(?:'s|'s|'s)\s+response){_SEPARATOR}(.+)",
        re.DOTALL | re.IGNORECASE,
    )

    followup_prefix = ""
    question_body = request.content.strip()

    m = _RE_ENHANCE_STAGE.match(question_body) or _RE_ENHANCE_MODEL.match(question_body)
    if m:
        followup_prefix = m.group(1).strip()   # e.g. "Regarding Stage 2"
        question_body = m.group(2).strip()      # the actual question
        logger.info("[enhance-prompt] Detected follow-up prefix: '%s'", followup_prefix)

    # ── System prompt for the enhancer ───────────────────────────
    enhance_system = """You are an expert prompt engineer for a pharmaceutical research LLM council. Your job is to MODESTLY improve a user's prompt — not to inflate or fabricate.

CRITICAL RULES:
1. NEVER invent context that is not in the original prompt. If the user asks about something you do not recognise, do NOT assume it is a drug, protein, algorithm, or anything else — keep the question as-is and only add minor clarifying structure.
2. Keep enhanced prompts PROPORTIONAL to the original length. A one-sentence question should become at most 2-3 sentences. Never turn a 5-word question into a paragraph.
3. If the topic is clearly pharmaceutical/scientific, you may add relevant dimensions (mechanism of action, safety profile, clinical evidence, regulatory status) — but only dimensions that actually apply.
4. If the topic is unknown, ambiguous, or not scientific, improve ONLY clarity and specificity. Do not force it into a scientific frame.
5. PRESERVE the user's actual intent and wording. Do not replace their terminology with synonyms or generalisations.
6. Do NOT add boilerplate phrases like "Provide a comprehensive scientific and technical elucidation" or "Elaborate on its utility across stages such as…". Write naturally.
7. Do NOT answer the question — only improve the phrasing.
8. Return ONLY the improved prompt text — no preamble, no explanation, no surrounding quotes.
9. Do NOT add any follow-up routing prefix (like "Regarding Stage 2:" or "About gpt-5.2's response:"). The system handles routing separately — you should only return the enhanced QUESTION.

EXAMPLES:
- Input: "What is metformin?" → "What is metformin, including its mechanism of action, primary indications, and key safety considerations?"
- Input: "What is clawdbot?" → "What is clawdbot?"  (unknown term — return as-is or nearly as-is)
- Input: "Compare SGLT2 inhibitors" → "Compare the major SGLT2 inhibitors (empagliflozin, dapagliflozin, canagliflozin) in terms of cardiovascular outcomes, renal benefits, and safety profiles based on recent clinical trial data."
- Input: "Tell me a joke" → "Tell me a joke"  (off-topic — return as-is)
- Input: "Need bit more clarity on FN:3" → "Could you elaborate on FN:3 — Under-emphasizes emicizumab's role — with more detail on the supporting evidence?"  (follow-up question — improve clarity, keep original references)"""

    messages = [
        {"role": "system", "content": enhance_system},
        {"role": "user", "content": f"Enhance this prompt (follow the critical rules strictly):\n\n{question_body}"}
    ]

    try:
        response = await query_model("gemini-2.5-flash", messages, timeout=30.0)
        if response and response.get('content'):
            enhanced_body = response['content'].strip().strip('"\'')

            # Re-attach the follow-up routing prefix if one was detected
            if followup_prefix:
                enhanced = f"{followup_prefix}: {enhanced_body}"
                logger.info("[enhance-prompt] Re-attached prefix → '%s: %s…'",
                            followup_prefix, enhanced_body[:60])
            else:
                enhanced = enhanced_body

            return {
                "original": request.content,
                "enhanced": enhanced,
            }
        else:
            raise HTTPException(status_code=502, detail="Failed to generate enhanced prompt")
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error enhancing prompt: {e}")
        raise HTTPException(status_code=502, detail="Failed to enhance prompt")


@app.get("/api/conversations", response_model=List[ConversationMetadata])
async def list_conversations(user_id: str = Depends(get_authenticated_user_id)):
    """List all conversations for the authenticated user (metadata only)."""
    convs = storage.list_conversations(user_id)

    # ── Auto-backfill missing context_tags ───────────────────────
    # Conversations created before context classification was introduced
    # will have no domain label in the sidebar.  Backfill them lazily
    # on each list request (lightweight regex classifier, no LLM call).
    backfilled = False
    for conv in convs:
        if conv.get("context_tags") or conv.get("message_count", 0) == 0:
            continue
        try:
            full = storage.get_conversation(user_id, conv["id"])
            if not full:
                continue
            # Find the first user message for classification
            first_user_msg = ""
            for msg in full.get("messages", []):
                if msg.get("role") == "user":
                    first_user_msg = msg.get("content", "")
                    break
            if not first_user_msg:
                continue
            from .memory import UserProfileMemory
            tags = UserProfileMemory.classify_query(first_user_msg)
            storage.update_conversation_context(user_id, conv["id"], tags)
            conv["context_tags"] = tags
            backfilled = True
        except Exception as e:
            logger.debug(f"[ContextBackfill] Non-fatal for {conv['id']}: {e}")

    if backfilled:
        logger.info(f"[ContextBackfill] Backfilled context_tags for user {user_id}")

    return convs


@app.post("/api/conversations", response_model=Conversation)
async def create_conversation(request: CreateConversationRequest, user_id: str = Depends(get_authenticated_user_id)):
    """Create a new conversation."""
    conversation_id = str(uuid.uuid4())
    conversation = storage.create_conversation(user_id, conversation_id)
    return conversation


@app.get("/api/conversations/{conversation_id}", response_model=Conversation)
async def get_conversation(conversation_id: str, user_id: str = Depends(get_authenticated_user_id)):
    """Get a specific conversation with all its messages."""
    conversation = storage.get_conversation(user_id, conversation_id)
    if conversation is None:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conversation


@app.get("/api/conversations/{conversation_id}/export")
async def export_conversation(conversation_id: str, user_id: str = Depends(get_authenticated_user_id), format: str = "markdown"):
    """
    Export a conversation in the specified format.
    
    Args:
        conversation_id: The conversation ID
        format: Export format - 'markdown', 'json', 'docx', or 'pptx' (default: markdown)
    
    Returns:
        markdown/json → JSON with {filename, content, content_type}
        docx/pptx     → Binary file download (StreamingResponse)
    """
    conversation = storage.get_conversation(user_id, conversation_id)
    if conversation is None:
        raise HTTPException(status_code=404, detail="Conversation not found")
    
    safe_title = re.sub(r'[^\w\s-]', '', conversation.get('title', 'conversation')).strip()[:80] or 'conversation'

    # ── DOCX export ──────────────────────────────────────────────
    if format == "docx":
        from .export_docx import generate_docx
        docx_bytes = generate_docx(conversation)
        return StreamingResponse(
            io.BytesIO(docx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{safe_title}.docx"'},
        )

    # ── PPTX export ──────────────────────────────────────────────
    if format == "pptx":
        from .export_pptx import generate_pptx
        pptx_bytes = generate_pptx(conversation)
        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{safe_title}.pptx"'},
        )

    # ── JSON export ──────────────────────────────────────────────
    if format == "json":
        return {
            "filename": f"{safe_title}.json",
            "content": json.dumps(conversation, indent=2),
            "content_type": "application/json"
        }
    
    # ── Default: Markdown export ─────────────────────────────────
    md_lines = [
        f"# {conversation['title']}",
        f"",
        f"**Created:** {conversation['created_at']}",
        f"",
        "---",
        ""
    ]
    
    for msg in conversation["messages"]:
        if msg["role"] == "user":
            md_lines.append(f"## 🧑 User")
            md_lines.append(f"")
            md_lines.append(msg["content"])
            md_lines.append("")
        else:
            # Assistant message with stages
            md_lines.append("## 🤖 Council Response")
            md_lines.append("")
            
            # Stage 1 - Individual Responses
            if "stage1" in msg and msg["stage1"]:
                md_lines.append("### Stage 1: Individual Model Responses")
                md_lines.append("")
                for response in msg["stage1"]:
                    md_lines.append(f"**{response['model']}:**")
                    md_lines.append(f"")
                    md_lines.append(response.get("response", "No response"))
                    md_lines.append("")
            
            # Stage 2 - Rankings
            if "stage2" in msg and msg["stage2"]:
                md_lines.append("### Stage 2: Peer Rankings")
                md_lines.append("")
                for ranking in msg["stage2"]:
                    md_lines.append(f"**{ranking['model']}:**")
                    md_lines.append(f"")
                    md_lines.append(ranking.get("response", "No ranking"))
                    md_lines.append("")
            
            # Stage 3 - Final Synthesis
            if "stage3" in msg and msg["stage3"]:
                md_lines.append("### Stage 3: Chairman's Final Synthesis")
                md_lines.append("")
                md_lines.append(f"**{msg['stage3'].get('model', 'Chairman')}:**")
                md_lines.append("")
                md_lines.append(msg["stage3"].get("response", "No synthesis"))
                md_lines.append("")
            
            md_lines.append("---")
            md_lines.append("")
    
    return {
        "filename": f"{safe_title}.md",
        "content": "\n".join(md_lines),
        "content_type": "text/markdown"
    }


@app.post("/api/conversations/{conversation_id}/message")
async def send_message(conversation_id: str, request: SendMessageRequest, user_id: str = Depends(get_authenticated_user_id)):
    """
    Send a message and run the 3-stage council process.
    Returns the complete response with all stages.
    """
    # Scope memory to this user
    set_memory_user(user_id)
    # Check if conversation exists
    conversation = storage.get_conversation(user_id, conversation_id)
    if conversation is None:
        raise HTTPException(status_code=404, detail="Conversation not found")

    # Check if this is the first message
    is_first_message = len(conversation["messages"]) == 0

    # Add user message
    storage.add_user_message(user_id, conversation_id, request.content)

    # If this is the first message, generate a title
    if is_first_message:
        title = await generate_conversation_title(request.content)
        storage.update_conversation_title(user_id, conversation_id, title)

    # Run the 3-stage council process with user preferences
    stage1_results, stage2_results, stage3_result, metadata = await run_full_council(
        request.content,
        council_models=request.council_models,
        chairman_model=request.chairman_model
    )

    # Add assistant message with all stages
    storage.add_assistant_message(
        user_id,
        conversation_id,
        stage1_results,
        stage2_results,
        stage3_result
    )

    # Return the complete response with metadata
    return {
        "stage1": stage1_results,
        "stage2": stage2_results,
        "stage3": stage3_result,
        "metadata": metadata
    }


@app.delete("/api/conversations/{conversation_id}")
async def delete_conversation(conversation_id: str, user_id: str = Depends(get_authenticated_user_id)):
    """Delete a conversation."""
    success = storage.delete_conversation(user_id, conversation_id)
    if not success:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return {"status": "deleted"}


@app.post("/api/conversations/{conversation_id}/analyze-agents")
async def analyze_agents(conversation_id: str, user_id: str = Depends(get_authenticated_user_id)):
    """Run agent team analysis on-demand for an existing conversation.

    Useful for conversations that were created before agent-team
    persistence was added, or when the original analysis failed.
    """
    set_memory_user(user_id)
    conv = storage.get_conversation(user_id, conversation_id)
    if conv is None:
        raise HTTPException(status_code=404, detail="Conversation not found")

    # Find the last assistant message
    msgs = conv.get("messages", [])
    last_assistant = None
    for msg in reversed(msgs):
        if msg.get("role") == "assistant":
            last_assistant = msg
            break

    if not last_assistant:
        raise HTTPException(status_code=400, detail="No assistant message found")

    stage1 = last_assistant.get("stage1", [])
    stage2 = last_assistant.get("stage2", [])
    stage3 = last_assistant.get("stage3", {})
    meta = last_assistant.get("metadata", {})

    if not stage1 or not stage3:
        raise HTTPException(status_code=400, detail="Incomplete conversation data")

    # Find the user query
    last_user = None
    for msg in reversed(msgs):
        if msg.get("role") == "user":
            last_user = msg
            break

    user_query = last_user.get("content", "") if last_user else ""

    # Run agent team
    try:
        result = await run_agent_team(
            user_query=user_query,
            stage1_results=stage1,
            stage2_results=stage2,
            stage3_result=stage3,
            aggregate_rankings=meta.get("aggregate_rankings", []),
            grounding_scores=meta.get("grounding_scores", {}),
            evidence_bundle=meta.get("evidence"),
            cost_summary=None,
        )
    except Exception as e:
        logger.error(f"On-demand agent analysis failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

    # Persist to storage
    try:
        storage.update_last_message_metadata(
            user_id, conversation_id, {"agent_team": result}
        )
    except Exception as e:
        logger.warning(f"Failed to persist on-demand agent_team: {e}")

    return result


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Kill Switch & Health Monitoring API
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


class KillSessionRequest(BaseModel):
    """Request to kill a specific session."""
    session_id: str
    reason: str = "User triggered kill switch"


class GlobalHaltRequest(BaseModel):
    """Request to activate global emergency halt."""
    reason: str = "Emergency halt triggered by user"


@app.post("/api/kill-switch/session")
async def kill_session(request: KillSessionRequest):
    """
    Kill Switch — abort a specific in-flight council session.
    This is the PRIMARY user-facing kill switch.
    """
    killed = kill_switch.kill_session(request.session_id, request.reason)
    if killed:
        return {
            "status": "killed",
            "session_id": request.session_id,
            "reason": request.reason,
        }
    raise HTTPException(
        status_code=404,
        detail=f"Session {request.session_id} not found or already completed",
    )


@app.post("/api/kill-switch/halt")
async def global_halt(request: GlobalHaltRequest):
    """
    Emergency Global Halt — kill ALL active sessions and block new ones.
    Use only in emergencies.
    """
    kill_switch.global_halt(request.reason)
    return {
        "status": "halted",
        "reason": request.reason,
        "sessions_killed": kill_switch.status()["active_session_count"],
    }


@app.post("/api/kill-switch/release")
async def release_halt():
    """Release global halt so new sessions can proceed."""
    kill_switch.release_global_halt()
    return {"status": "released"}


@app.get("/api/kill-switch/status")
async def get_kill_switch_status():
    """Get current kill switch status (active sessions, halt state)."""
    return kill_switch.status()


@app.get("/api/health")
async def get_system_health():
    """
    Full system health: kill switch state, circuit breaker per-model status,
    recent self-healing actions taken, and security configuration.
    """
    status = health_monitor.full_status()
    status["security"] = get_security_status()
    return status


@app.get("/api/health/circuits")
async def get_circuit_status():
    """Get per-model circuit breaker status."""
    return circuit_breaker.status()


@app.post("/api/health/circuits/reset")
async def reset_circuits(model: Optional[str] = None):
    """Reset circuit breaker for a specific model or all models."""
    circuit_breaker.reset(model)
    return {
        "status": "reset",
        "model": model or "all",
    }


# ────────────────────────────────────────────────────────────────────────
# A2A Agent Card Discovery (/.well-known/agent-card.json)
# ────────────────────────────────────────────────────────────────────────

_WELLKNOWN_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), ".well-known")


def _load_json(path: str):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


@app.get("/.well-known/agent-card.json")
async def get_agent_card():
    """A2A protocol discovery endpoint — returns the main council agent card."""
    card_path = os.path.join(_WELLKNOWN_DIR, "agent-card.json")
    if not os.path.exists(card_path):
        raise HTTPException(status_code=404, detail="Agent card not found")
    return _load_json(card_path)


@app.get("/api/agent-cards")
async def list_agent_cards():
    """List all individual agent cards (core + VP)."""
    agents_dir = os.path.join(_WELLKNOWN_DIR, "agents")
    if not os.path.isdir(agents_dir):
        return {"agents": []}
    cards = []
    for fname in sorted(os.listdir(agents_dir)):
        if fname.endswith(".json"):
            cards.append(_load_json(os.path.join(agents_dir, fname)))
    return {"agents": cards, "count": len(cards)}


@app.get("/api/agent-cards/{agent_id}")
async def get_individual_agent_card(agent_id: str):
    """Get a specific agent card by ID (e.g. 'research-analyst')."""
    card_path = os.path.join(_WELLKNOWN_DIR, "agents", f"{agent_id}.json")
    if not os.path.exists(card_path):
        raise HTTPException(status_code=404, detail=f"Agent card '{agent_id}' not found")
    return _load_json(card_path)


@app.get("/api/agent-cards-download")
async def download_agent_cards():
    """Download the full A2A agent card bundle as a single JSON file."""
    from fastapi.responses import Response

    # Load the main council card
    main_card_path = os.path.join(_WELLKNOWN_DIR, "agent-card.json")
    main_card = _load_json(main_card_path) if os.path.exists(main_card_path) else {}

    # Load all individual agent cards
    agents_dir = os.path.join(_WELLKNOWN_DIR, "agents")
    agent_cards = []
    if os.path.isdir(agents_dir):
        for fname in sorted(os.listdir(agents_dir)):
            if fname.endswith(".json"):
                agent_cards.append(_load_json(os.path.join(agents_dir, fname)))

    bundle = {
        "a2a_protocol_version": "1.0",
        "generated_at": datetime.utcnow().isoformat() + "Z",
        "council": main_card,
        "agents": agent_cards,
        "agent_count": len(agent_cards),
    }

    content = json.dumps(bundle, indent=2, ensure_ascii=False)
    return Response(
        content=content,
        media_type="application/json",
        headers={
            "Content-Disposition": 'attachment; filename="llm-council-agent-cards.json"'
        },
    )


# ────────────────────────────────────────────────────────────────────────
# Memory Management API
# ────────────────────────────────────────────────────────────────────────

class MemoryDecisionRequest(BaseModel):
    decision: str  # "learn" | "unlearn"
    memory_type: str  # "semantic" | "episodic" | "procedural"
    memory_id: str
    reason: Optional[str] = ""


@app.get("/api/memory/stats")
async def get_memory_stats(user_id: str = Depends(get_authenticated_user_id)):
    """Get memory statistics across all three tiers for the current user."""
    set_memory_user(user_id)
    mm = get_memory_manager()
    return mm.stats()


@app.get("/api/memory/{memory_type}")
async def list_memories(memory_type: str, include_unlearned: bool = False, user_id: str = Depends(get_authenticated_user_id)):
    """List all memories for a given type (semantic, episodic, procedural), scoped to user."""
    set_memory_user(user_id)
    mm = get_memory_manager()
    if memory_type == "semantic":
        return mm.semantic.list_all(include_unlearned=include_unlearned)
    elif memory_type == "episodic":
        return mm.episodic.list_all(include_unlearned=include_unlearned)
    elif memory_type == "procedural":
        return mm.procedural.list_all(include_unlearned=include_unlearned)
    else:
        raise HTTPException(status_code=400, detail=f"Unknown memory type: {memory_type}")


@app.get("/api/memory/{memory_type}/{memory_id}")
async def get_memory_entry(memory_type: str, memory_id: str, user_id: str = Depends(get_authenticated_user_id)):
    """Get a specific memory entry, scoped to user."""
    set_memory_user(user_id)
    from .memory_store import get_memory_backend
    backend = get_memory_backend()
    doc = backend.get(memory_type, memory_id)
    if not doc:
        raise HTTPException(status_code=404, detail="Memory entry not found")
    return doc


@app.post("/api/memory/decision")
async def apply_memory_decision(request: MemoryDecisionRequest, user_id: str = Depends(get_authenticated_user_id)):
    """Apply a learn/unlearn decision from the user."""
    set_memory_user(user_id)
    result = await user_gate_agent(
        decision=request.decision,
        memory_type=request.memory_type,
        memory_id=request.memory_id,
        reason=request.reason or "",
        user_id=user_id,
    )
    if not result.get("success"):
        raise HTTPException(status_code=404, detail="Memory entry not found or invalid decision")
    return result


@app.get("/api/memory/search/{memory_type}")
async def search_memories(memory_type: str, q: str, limit: int = 10, user_id: str = Depends(get_authenticated_user_id)):
    """Search memories by text query, scoped to user."""
    set_memory_user(user_id)
    from .memory_store import get_memory_backend
    backend = get_memory_backend()
    return backend.search(memory_type, q, limit=limit)


@app.delete("/api/memory/{memory_type}/{memory_id}")
async def delete_memory_entry(memory_type: str, memory_id: str, user_id: str = Depends(get_authenticated_user_id)):
    """Permanently delete a memory entry, scoped to user."""
    set_memory_user(user_id)
    from .memory_store import get_memory_backend
    backend = get_memory_backend()
    if backend.delete(memory_type, memory_id):
        return {"status": "deleted", "memory_type": memory_type, "memory_id": memory_id}
    raise HTTPException(status_code=404, detail="Memory entry not found")


# ═══════════════════════════════════════════════════════════════════════
#  Targeted Follow-Up Helpers
#  When the user clicks a FOCUS ON chip, we skip the 3-stage pipeline
#  and route directly to the chairman with the referenced content.
# ═══════════════════════════════════════════════════════════════════════

# ── Targeted Follow-Up Patterns ────────────────────────────────────
# These detect when a user is asking about a SPECIFIC stage or model
# from the previous council response.  Detection is intentionally
# generous — both chip-generated prefixes ("Regarding Stage 3: …")
# and free-form user text ("Regarding Stage 3, …" or "About Stage
# 3 - tell me more") should match.
#
# Separator after "Stage N" / "response" is flexible:
#   colon, comma, semicolon, dash (any kind), period, or nothing.
#
# Prefixes accepted: Regarding, About, Re, On, For, Expand on,
#   Elaborate on, Tell me more about, More on, Concerning.
# ──────────────────────────────────────────────────────────────────

_FOLLOWUP_PREFIX = r"(?:Regarding|About|Re:?|On|For|Expand\s+on|Elaborate\s+on|Tell\s+me\s+more\s+about|More\s+on|Concerning)"
_SEPARATOR = r"\s*[,:;\-–—.!?]?\s*"

_RE_STAGE_FOLLOWUP = re.compile(
    rf"^{_FOLLOWUP_PREFIX}\s+(Stage\s*\d){_SEPARATOR}(.+)",
    re.DOTALL | re.IGNORECASE,
)
_RE_MODEL_FOLLOWUP = re.compile(
    rf"^{_FOLLOWUP_PREFIX}\s+(.+?)(?:'s|'s|'s)\s+response{_SEPARATOR}(.+)",
    re.DOTALL | re.IGNORECASE,
)

# Fallback: detect "Stage N" anywhere in the first 60 chars when there
# IS a previous assistant message — catches patterns like
# "please list … from Stage 3" or "what did Stage 1 say about…"
_RE_STAGE_MENTION = re.compile(
    r"\b(Stage\s*(\d))\b", re.IGNORECASE
)

# Fallback: detect a known model short-name anywhere in the first 80
# chars — catches "what did gpt-5.2 think about …"
# (populated dynamically from the previous assistant message's Stage 1)


def _detect_targeted_followup(
    content: str, conversation_history: list
) -> Optional[dict]:
    """Return a dict describing the targeted follow-up, or *None*.

    The returned dict has keys:
        type          – "stage" | "model"
        target_label  – e.g. "Stage 3" or "gpt-5.2"
        user_question – the actual question text (after the prefix)
        reference_data – the referenced stage/model content from the last
                         assistant message
        prev_stage1   – full Stage 1 results (for agent team)
        prev_stage2   – full Stage 2 results (for agent team)
        prev_stage3   – full Stage 3 result  (for agent team)
        prev_metadata – metadata from last assistant message
    """
    if not content or not conversation_history:
        return None

    # Find the last assistant message
    last_assistant = None
    for msg in reversed(conversation_history):
        if msg.get("role") == "assistant":
            last_assistant = msg
            break
    if last_assistant is None:
        return None

    def _build_stage_result(stage_label: str, user_question: str) -> dict:
        """Build the targeted follow-up dict for a stage reference."""
        stage_num = stage_label.strip()[-1]  # "1", "2", "3"
        reference_data = None
        if stage_num == "1":
            reference_data = last_assistant.get("stage1", [])
        elif stage_num == "2":
            reference_data = last_assistant.get("stage2", [])
        elif stage_num == "3":
            s3 = last_assistant.get("stage3", {})
            reference_data = s3.get("response", "") if isinstance(s3, dict) else s3
        return {
            "type": "stage",
            "target_label": f"Stage {stage_num}",
            "user_question": user_question,
            "reference_data": reference_data,
            "prev_stage1": last_assistant.get("stage1", []),
            "prev_stage2": last_assistant.get("stage2", []),
            "prev_stage3": last_assistant.get("stage3", {}),
            "prev_metadata": last_assistant.get("metadata", {}),
        }

    def _build_model_result(model_short: str, full_model_id: str,
                            model_response: str, user_question: str) -> dict:
        """Build the targeted follow-up dict for a model reference."""
        return {
            "type": "model",
            "target_label": model_short,
            "full_model_id": full_model_id,
            "user_question": user_question,
            "reference_data": model_response,
            "prev_stage1": last_assistant.get("stage1", []),
            "prev_stage2": last_assistant.get("stage2", []),
            "prev_stage3": last_assistant.get("stage3", {}),
            "prev_metadata": last_assistant.get("metadata", {}),
        }

    # ── PRIMARY: Stage-targeted follow-up (prefix match) ──
    m = _RE_STAGE_FOLLOWUP.match(content)
    if m:
        stage_label = m.group(1)          # "Stage 1" / "Stage 2" / "Stage 3"
        user_question = m.group(2).strip()
        logger.info(f"[Targeted Follow-Up Detection] PRIMARY stage match: '{stage_label}' from prefix pattern")
        return _build_stage_result(stage_label, user_question)

    # ── PRIMARY: Model-targeted follow-up (prefix match) ──
    m = _RE_MODEL_FOLLOWUP.match(content)
    if m:
        model_short = m.group(1).strip()   # e.g. "gpt-5.2"
        user_question = m.group(2).strip()

        # Find the matching model's Stage 1 response
        s1_results = last_assistant.get("stage1", [])
        model_response = None
        full_model_id = model_short
        for r in s1_results:
            mid = r.get("model", "")
            short = mid.split("/")[-1] if "/" in mid else mid
            if short.lower() == model_short.lower() or mid.lower() == model_short.lower():
                model_response = r.get("response", "")
                full_model_id = mid
                break

        if model_response is not None:
            logger.info(f"[Targeted Follow-Up Detection] PRIMARY model match: '{model_short}' from prefix pattern")
            return _build_model_result(model_short, full_model_id, model_response, user_question)

    # ── FALLBACK: Stage mention in the first 80 chars ──
    # Catches: "please list hemophilia centers from Stage 3",
    #          "can you expand the Stage 1 answer", etc.
    head = content[:80]
    stage_mentions = list(_RE_STAGE_MENTION.finditer(head))
    if stage_mentions:
        # Use the LAST Stage mention (most likely the target)
        last_match = stage_mentions[-1]
        stage_label = last_match.group(1)  # e.g. "Stage 3"
        # Use full content as the question (can't reliably split)
        logger.info(f"[Targeted Follow-Up Detection] FALLBACK stage mention: '{stage_label}' in first 80 chars")
        return _build_stage_result(stage_label, content)

    # ── FALLBACK: Model short-name mention in the first 80 chars ──
    s1_results = last_assistant.get("stage1", [])
    if s1_results:
        content_lower = head.lower()
        for r in s1_results:
            mid = r.get("model", "")
            short = (mid.split("/")[-1] if "/" in mid else mid).lower()
            if short and short in content_lower:
                model_response = r.get("response", "")
                logger.info(f"[Targeted Follow-Up Detection] FALLBACK model mention: '{short}' in first 80 chars")
                return _build_model_result(short, mid, model_response, content)

    logger.debug(f"[Targeted Follow-Up Detection] No match for: {content[:100]!r}")
    return None


async def _run_targeted_followup(
    *,
    targeted: dict,
    user_query: str,
    conversation_history: list,
    conversation_id: str,
    user_id: str,
    user_chairman_model: Optional[str],
    user_council_models: Optional[list],
    web_search_enabled: bool,
    speed_mode: bool,
    session_id: str,
    SPEED_TIMEOUT: float,
    SPEED_S3_MAX_TOKENS: Optional[int],
    cost_tracker,
):
    """Async generator yielding SSE events for a targeted follow-up.

    Routing rules:
    - Stage 1 target  → Re-query all council models with the follow-up question
    - Stage 2 target  → Re-query all council models for peer evaluation of the follow-up
    - Stage 3 target  → Chairman-only synthesis (fast path)
    - Model target    → Chairman answers about that specific model's response (fast path)

    The agent team still fires for analysis in all cases.
    """
    from .config import DEFAULT_CHAIRMAN_MODEL as _DEFAULT_CHAIRMAN
    from .config import COUNCIL_MODELS as _COUNCIL_MODELS

    chairman = user_chairman_model or _DEFAULT_CHAIRMAN
    target_type = targeted["type"]
    target_label = targeted["target_label"]
    user_question = targeted["user_question"]
    reference_data = targeted["reference_data"]

    # Re-use previous pipeline results for agents
    prev_s1 = targeted["prev_stage1"]
    prev_s2 = targeted["prev_stage2"]
    prev_s3 = targeted["prev_stage3"]
    prev_meta = targeted["prev_metadata"]
    prev_rankings = prev_meta.get("aggregate_rankings", [])
    prev_grounding = prev_meta.get("grounding_scores", {})
    prev_evidence = prev_meta.get("evidence", {})

    # ── Tell frontend we are on the targeted path ──
    yield f"data: {json.dumps({'type': 'targeted_followup_start', 'data': {'target_type': target_type, 'target': target_label}})}\n\n"

    # ── Classify the follow-up query for context tagging ──
    tf_context_tags = None
    try:
        from .memory import UserProfileMemory
        tf_context_tags = UserProfileMemory.classify_query(user_query)
        storage.update_conversation_context(user_id, conversation_id, tf_context_tags)
        yield f"data: {json.dumps({'type': 'context_classified', 'data': tf_context_tags})}\n\n"
    except Exception as e:
        logger.warning(f"[TargetedFU-ContextTags] Non-fatal: {e}")

    # ── Route based on target ────────────────────────────────────
    if target_type == "stage" and target_label == "Stage 1":
        # ━━ STAGE 1 FOLLOW-UP: Re-query all council models ━━━━━━
        logger.info("[Targeted Follow-Up] Stage 1 → re-running council models")
        yield f"data: {json.dumps({'type': 'stage1_start'})}\n\n"

        s1_models = user_council_models or _COUNCIL_MODELS
        s1_context = build_conversation_context(conversation_history)
        s1_query = f"{s1_context}Current question (follow-up): {user_question}" if s1_context else user_question
        s1_messages = [{"role": "user", "content": s1_query}]

        SPEED_S1_MAX_TOKENS = 2048 if speed_mode else None

        pending_tasks = {}
        for model in s1_models:
            task = asyncio.create_task(
                query_model(model, s1_messages,
                            timeout=SPEED_TIMEOUT,
                            web_search_enabled=web_search_enabled,
                            session_id=session_id,
                            max_retries=1,
                            max_tokens=SPEED_S1_MAX_TOKENS)
            )
            pending_tasks[task] = model

        stage1_results = []
        s1_failed_models = []
        s1_used_models = set(s1_models)
        s1_total = len(s1_models)

        while pending_tasks:
            done, _ = await asyncio.wait(
                pending_tasks.keys(), return_when=asyncio.FIRST_COMPLETED
            )
            for task in done:
                model = pending_tasks.pop(task)
                try:
                    response = task.result()
                    if response is not None:
                        result_item = {
                            "model": model,
                            "response": response.get("content", ""),
                            "usage": response.get("usage"),
                        }
                        stage1_results.append(result_item)
                        yield f"data: {json.dumps({'type': 'stage1_model_complete', 'data': result_item, 'progress': {'completed': len(stage1_results), 'failed': len(s1_failed_models), 'total': s1_total}})}\n\n"
                    else:
                        s1_failed_models.append(model)
                except Exception as e:
                    logger.error(f"[Targeted S1] {model} raised: {e}")
                    s1_failed_models.append(model)

        # Self-healing fallbacks
        if s1_failed_models and len(stage1_results) < len(s1_models):
            for failed_model in s1_failed_models:
                fallback = resolve_fallback(failed_model, s1_used_models)
                if fallback:
                    s1_used_models.add(fallback)
                    fb_resp = await query_model(fallback, s1_messages,
                                                web_search_enabled=web_search_enabled,
                                                session_id=session_id)
                    if fb_resp is not None:
                        result_item = {
                            "model": f"{fallback} (fallback for {failed_model})",
                            "response": fb_resp.get("content", ""),
                            "usage": fb_resp.get("usage"),
                        }
                        stage1_results.append(result_item)
                        yield f"data: {json.dumps({'type': 'stage1_model_complete', 'data': result_item, 'progress': {'completed': len(stage1_results), 'failed': len(s1_failed_models), 'total': s1_total}})}\n\n"

        for r in stage1_results:
            cost_tracker.record("stage1", r["model"], r.get("usage"))
        yield f"data: {json.dumps({'type': 'stage1_complete', 'data': stage1_results})}\n\n"

        # Save — Stage 1 only, carry forward previous Stage 2/3
        storage.add_assistant_message(
            user_id, conversation_id,
            stage1_results, prev_s2, prev_s3,
            metadata={
                "label_to_model": prev_meta.get("label_to_model", {}),
                "aggregate_rankings": prev_rankings,
                "grounding_scores": prev_grounding,
                "evidence": prev_evidence,
                "targeted_followup": {"type": target_type, "target": target_label},
                "context_tags": tf_context_tags,
            },
        )

        # Agent team
        agent_team_result = None
        try:
            agent_team_result = await run_agent_team(
                user_query=user_query,
                stage1_results=stage1_results,
                stage2_results=prev_s2,
                stage3_result=prev_s3,
                aggregate_rankings=prev_rankings,
                grounding_scores=prev_grounding,
                evidence_bundle=prev_evidence,
                cost_summary=cost_tracker.summary(),
                web_search_enabled=web_search_enabled,
            )
        except Exception as e:
            logger.warning(f"[Targeted S1] Agent team error: {e}")

        cost_summary = cost_tracker.summary()
        yield f"data: {json.dumps({'type': 'cost_summary', 'data': cost_summary})}\n\n"
        if agent_team_result:
            yield f"data: {json.dumps({'type': 'agent_team_complete', 'data': agent_team_result})}\n\n"

        yield f"data: {json.dumps({'type': 'complete'})}\n\n"
        return

    elif target_type == "stage" and target_label == "Stage 2":
        # ━━ STAGE 2 FOLLOW-UP: Re-run actual peer evaluations ━━━━━━━
        # Uses all council models to re-evaluate Stage 1 responses,
        # injecting the follow-up question as additional context.
        # Emits stage2_start → stage2_model_response → stage2_complete
        # (NOT stage3 events — Stage 2 stays as Stage 2).
        logger.info("[Targeted Follow-Up] Stage 2 → re-running peer evaluations with all council models")

        yield f"data: {json.dumps({'type': 'stage2_start'})}\n\n"

        # ── Build Stage 2 prompt with follow-up context ──
        SPEED_S2_MAX_TOKENS = 1536 if speed_mode else None
        s2_models = user_council_models or _COUNCIL_MODELS
        s2_labels = [chr(65 + i) for i in range(len(prev_s1 or []))]
        label_to_model = {
            f"Response {label}": result['model']
            for label, result in zip(s2_labels, prev_s1 or [])
        }

        # Build conversation context for follow-up note
        s2_context = build_conversation_context(conversation_history)
        s2_context_note = f"\nNote: This is a follow-up question in an ongoing conversation.\n{s2_context}\n" if s2_context else ""

        # Augmented question includes follow-up context
        augmented_question = f"""{user_query}

[FOLLOW-UP CONTEXT] The user has a specific follow-up regarding the peer evaluations:
{user_question}

Please re-evaluate the responses with this follow-up concern in mind."""

        # ── Position Debiasing per reviewer ──
        s2_per_model_responses = {}
        for m in s2_models:
            shuf = list(range(len(prev_s1 or [])))
            import random
            random.shuffle(shuf)
            s2_per_model_responses[m] = "\n\n".join([
                f"Response {s2_labels[i]}:\n{(prev_s1 or [])[i].get('response', '')}"
                for i in shuf
            ])

        # Use the same prompt template as the main pipeline
        s2_prompt_template = """You are a pharmaceutical domain expert evaluating different responses to the following question:
{context_note}
Question: {question}

Here are the responses from different models (anonymized):

{responses_text}

═══════════════════════════════════════════════════════════
PART 1 — RUBRIC EVALUATION
═══════════════════════════════════════════════════════════
For EACH response, provide a score from 0 to 10 on each criterion below.
After each score, give a brief justification (1 sentence).

Criteria:
  • Relevancy (0-10): How directly and completely the response addresses the original question
  • Faithfulness (0-10): Factual accuracy, absence of hallucinations
  • Context Recall (0-10): Coverage of key concepts and nuances
  • Output Quality (0-10): Clarity, structure, depth
  • Consensus (0-10): Would other domain experts broadly agree?

Format EXACTLY as follows for each response:

RUBRIC Response X:
  Relevancy: <score>/10 — <justification>
  Faithfulness: <score>/10 — <justification>
  Context Recall: <score>/10 — <justification>
  Output Quality: <score>/10 — <justification>
  Consensus: <score>/10 — <justification>

═══════════════════════════════════════════════════════════
FINAL RANKING
═══════════════════════════════════════════════════════════
Based on your rubric evaluation above, provide your final ranking
from best to worst.

IMPORTANT: Your final ranking MUST be formatted EXACTLY as follows:
- Start with the line "FINAL RANKING:" (all caps, with colon)
- Then list the responses from best to worst as a numbered list
- Each line should be: number, period, space, then ONLY the response label (e.g., "1. Response A")

Now provide your complete evaluation:"""

        # Build per-model messages
        s2_per_model_msgs = {}
        for m in s2_models:
            s2_per_model_msgs[m] = [{"role": "user", "content": s2_prompt_template.format(
                context_note=s2_context_note,
                question=augmented_question,
                responses_text=s2_per_model_responses.get(m, ""),
            )}]

        # ── Fire all reviewers in parallel ──
        s2_tasks = {
            asyncio.create_task(
                query_model(m, s2_per_model_msgs[m],
                            timeout=SPEED_TIMEOUT,
                            web_search_enabled=web_search_enabled,
                            session_id=session_id,
                            max_tokens=SPEED_S2_MAX_TOKENS)
            ): m
            for m in s2_models
        }

        stage2_results = []
        s2_pending = set(s2_tasks.keys())
        s2_total = len(s2_models)

        while s2_pending:
            done, s2_pending = await asyncio.wait(s2_pending, return_when=asyncio.FIRST_COMPLETED)
            for task in done:
                model_name = s2_tasks[task]
                try:
                    resp = task.result()
                    if resp is not None:
                        full_text = resp.get('content', '')
                        parsed = parse_ranking_from_text(full_text)
                        rubric = parse_rubric_scores(full_text)
                        claims = parse_claim_counts(full_text)
                        result_item = {
                            "model": model_name,
                            "ranking": full_text,
                            "parsed_ranking": parsed,
                            "rubric_scores": rubric,
                            "claim_counts": claims,
                            "usage": resp.get('usage'),
                        }
                        stage2_results.append(result_item)
                        cost_tracker.record("stage2", model_name, resp.get('usage'))
                        yield f"data: {json.dumps({'type': 'stage2_model_response', 'data': result_item, 'progress': {'completed': len(stage2_results), 'total': s2_total}})}\n\n"
                except Exception as e:
                    logger.error(f"[Targeted S2] {model_name} failed: {e}")

        # Compute aggregate rankings and grounding scores
        aggregate_rankings = calculate_aggregate_rankings(stage2_results, label_to_model)
        grounding_scores = compute_response_grounding_scores(
            stage2_results, label_to_model, aggregate_rankings
        )

        yield f"data: {json.dumps({'type': 'stage2_complete', 'data': stage2_results, 'metadata': {'label_to_model': label_to_model, 'aggregate_rankings': aggregate_rankings, 'grounding_scores': grounding_scores}})}\n\n"

        cost_summary = cost_tracker.summary()
        yield f"data: {json.dumps({'type': 'cost_summary', 'data': cost_summary})}\n\n"

        agent_team_result = None
        try:
            agent_team_result = await run_agent_team(
                user_query=user_query, stage1_results=prev_s1,
                stage2_results=stage2_results, stage3_result=prev_s3,
                aggregate_rankings=aggregate_rankings, grounding_scores=grounding_scores,
                evidence_bundle=prev_evidence, cost_summary=cost_summary,
                web_search_enabled=web_search_enabled,
            )
        except Exception as e:
            logger.warning(f"[Targeted S2] Agent team error: {e}")
        if agent_team_result:
            yield f"data: {json.dumps({'type': 'agent_team_complete', 'data': agent_team_result})}\n\n"

        storage.add_assistant_message(
            user_id, conversation_id, prev_s1, stage2_results, prev_s3,
            metadata={
                "label_to_model": label_to_model,
                "aggregate_rankings": aggregate_rankings,
                "grounding_scores": grounding_scores,
                "evidence": prev_evidence,
                "targeted_followup": {"type": target_type, "target": target_label},
                "context_tags": tf_context_tags,
                **({"agent_team": agent_team_result} if agent_team_result else {}),
            },
        )
        yield f"data: {json.dumps({'type': 'complete'})}\n\n"
        return

    # ━━ MODEL TARGET: Direct model response ━━━━━━━━━━━━━━━━━━━━━
    # When the user targets a specific model, that MODEL answers directly
    # (not the chairman). Falls back to chairman if the model fails.
    if target_type == "model":
        target_model = targeted.get("full_model_id", target_label)
        logger.info("[Targeted Follow-Up] %s → direct model query to %s", target_label, target_model)

        direct_prompt = f"""You previously answered a question as part of a multi-LLM council deliberation.
The user is now asking YOU a follow-up question specifically about YOUR response.

--- YOUR PREVIOUS RESPONSE ---
{(reference_data or '')[:8000]}

--- USER'S FOLLOW-UP QUESTION ---
{user_question}

Answer the follow-up question directly and thoroughly in your own voice.
You are the council member being addressed — respond as the expert, not as a summariser.
If the user asks for elaboration, go deeper into the topic with additional detail.
If they ask for clarification, clarify your reasoning and provide supporting evidence.
Maintain pharmaceutical domain expertise and cite evidence where applicable."""

        yield f"data: {json.dumps({'type': 'stage3_start', 'data': {'direct_model': True, 'model': target_model}})}\n\n"

        direct_response = None
        try:
            direct_response = await query_model(
                target_model,
                [{"role": "user", "content": direct_prompt}],
                timeout=SPEED_TIMEOUT,
                max_tokens=SPEED_S3_MAX_TOKENS,
                session_id=session_id,
            )
        except Exception as e:
            logger.warning(f"[Targeted Follow-Up] Direct model {target_model} failed: {e}, falling back to chairman")

        # Fallback to chairman if target model fails
        if direct_response is None:
            logger.info(f"[Targeted Follow-Up] Falling back to chairman {chairman} for {target_label}")
            fallback_prompt = f"""The user asked a follow-up about {target_label}'s response, but that model
is unavailable. Please answer on behalf of the council.

--- REFERENCED RESPONSE ({target_label}) ---
{(reference_data or '')[:8000]}

--- FOLLOW-UP QUESTION ---
{user_question}

Provide a thorough, direct answer. Maintain pharmaceutical domain expertise."""
            try:
                direct_response = await query_model(
                    chairman,
                    [{"role": "user", "content": fallback_prompt}],
                    timeout=SPEED_TIMEOUT,
                    max_tokens=SPEED_S3_MAX_TOKENS,
                    session_id=session_id,
                )
                target_model = f"{chairman} (fallback for {target_label})"
            except Exception as e2:
                logger.error(f"[Targeted Follow-Up] Chairman fallback also failed: {e2}")
                yield f"data: {json.dumps({'type': 'error', 'message': f'Follow-up query failed: {e2}'})}\n\n"
                return

        stage3_result = {
            "model": target_model,
            "response": enrich_stage3_citations((direct_response or {}).get("content", "")),
            "direct_model": True,
        }

        yield f"data: {json.dumps({'type': 'stage3_complete', 'data': stage3_result})}\n\n"

        # Cost, agents, storage — same flow as below
        cost_summary = cost_tracker.summary()
        yield f"data: {json.dumps({'type': 'cost_summary', 'data': cost_summary})}\n\n"

        agent_team_result = None
        try:
            agent_team_result = await run_agent_team(
                user_query=user_query,
                stage1_results=prev_s1,
                stage2_results=prev_s2,
                stage3_result=stage3_result,
                aggregate_rankings=prev_rankings,
                grounding_scores=prev_grounding,
                evidence_bundle=prev_evidence,
                cost_summary=cost_summary,
                web_search_enabled=web_search_enabled,
            )
        except Exception as e:
            logger.warning(f"[Targeted Follow-Up] Agent team error: {e}")
        if agent_team_result:
            yield f"data: {json.dumps({'type': 'agent_team_complete', 'data': agent_team_result})}\n\n"

        storage.add_assistant_message(
            user_id, conversation_id, prev_s1, prev_s2, stage3_result,
            metadata={
                "label_to_model": prev_meta.get("label_to_model", {}),
                "aggregate_rankings": prev_rankings,
                "grounding_scores": prev_grounding,
                "evidence": prev_evidence,
                "targeted_followup": {"type": target_type, "target": target_label},
                "context_tags": tf_context_tags,
                **({"agent_team": agent_team_result} if agent_team_result else {}),
            },
        )
        yield f"data: {json.dumps({'type': 'complete'})}\n\n"
        return

    # ━━ STAGE 3 TARGET: Chairman fast path ━━━━━━━━━━━━━━━━━━━━━
    # Stage 3 follow-ups go to the chairman for re-synthesis
    logger.info("[Targeted Follow-Up] %s → chairman fast path", target_label)

    context_block = reference_data if isinstance(reference_data, str) else json.dumps(reference_data)

    focused_prompt = f"""The user previously asked a question and received a council deliberation.
Now they have a follow-up question specifically about {target_label}.

--- REFERENCED CONTENT ---
{context_block[:8000]}

--- FOLLOW-UP QUESTION ---
{user_question}

Provide a thorough, direct answer to the follow-up question.
Focus specifically on the referenced content above.
If the user asks for elaboration, provide deeper analysis.
If they ask for comparison, compare with what other council members said.
Maintain pharmaceutical domain expertise and cite evidence where applicable."""

    # ── Query chairman ──
    yield f"data: {json.dumps({'type': 'stage3_start'})}\n\n"

    try:
        chairman_response = await query_model(
            chairman,
            [{"role": "user", "content": focused_prompt}],
            timeout=SPEED_TIMEOUT,
            max_tokens=SPEED_S3_MAX_TOKENS,
            session_id=session_id,
        )
    except Exception as e:
        logger.error(f"[Targeted Follow-Up] Chairman query failed: {e}")
        yield f"data: {json.dumps({'type': 'error', 'message': f'Follow-up query failed: {e}'})}\n\n"
        return

    stage3_result = {
        "model": chairman,
        "response": (chairman_response or {}).get("content", ""),
    }

    # Enrich citations
    stage3_result["response"] = enrich_stage3_citations(stage3_result["response"])

    yield f"data: {json.dumps({'type': 'stage3_complete', 'data': stage3_result})}\n\n"

    # ── Cost tracking ──
    cost_summary = cost_tracker.summary()
    yield f"data: {json.dumps({'type': 'cost_summary', 'data': cost_summary})}\n\n"

    # ── Agent Team (same as full pipeline — still fires) ──
    agent_team_result = None
    try:
        agent_team_result = await run_agent_team(
            user_query=user_query,
            stage1_results=prev_s1,
            stage2_results=prev_s2,
            stage3_result=stage3_result,
            aggregate_rankings=prev_rankings,
            grounding_scores=prev_grounding,
            evidence_bundle=prev_evidence,
            cost_summary=cost_summary,
            web_search_enabled=web_search_enabled,
        )
    except Exception as e:
        logger.warning(f"[Targeted Follow-Up] Agent team error: {e}")

    if agent_team_result:
        yield f"data: {json.dumps({'type': 'agent_team_complete', 'data': agent_team_result})}\n\n"

    # ── Save assistant message ──
    storage.add_assistant_message(
        user_id,
        conversation_id,
        prev_s1,       # carry forward Stage 1
        prev_s2,       # carry forward Stage 2
        stage3_result,
        metadata={
            "label_to_model": prev_meta.get("label_to_model", {}),
            "aggregate_rankings": prev_rankings,
            "grounding_scores": prev_grounding,
            "evidence": prev_evidence,
            "targeted_followup": {
                "type": target_type,
                "target": target_label,
            },
            "context_tags": tf_context_tags,
            **({"agent_team": agent_team_result} if agent_team_result else {}),
        },
    )

    yield f"data: {json.dumps({'type': 'complete'})}\n\n"


@app.post("/api/conversations/{conversation_id}/message/stream")
async def send_message_stream(conversation_id: str, request: SendMessageRequest, user_id: str = Depends(get_authenticated_user_id)):
    """
    Send a message and stream the 3-stage council process.
    Returns Server-Sent Events as each stage completes.
    Supports file attachments (PDF, PPTX, XLSX, DOCX).
    """
    # Check if conversation exists
    conversation = storage.get_conversation(user_id, conversation_id)
    if conversation is None:
        raise HTTPException(status_code=404, detail="Conversation not found")

    # Validate attachments
    for attachment in request.attachments:
        error = validate_attachment(attachment)
        if error:
            raise HTTPException(status_code=400, detail=error)

    # Build augmented content with file attachments info
    augmented_content = request.content
    if request.attachments:
        attachment_descriptions = [
            extract_file_content_description(att) for att in request.attachments
        ]
        if augmented_content:
            augmented_content = f"{augmented_content}\n\n---\nAttached Files:\n" + "\n".join(attachment_descriptions)
        else:
            augmented_content = "Please analyze the following attached files:\n" + "\n".join(attachment_descriptions)

    # Check if this is the first message
    is_first_message = len(conversation["messages"]) == 0
    
    # Get user preferences for council/chairman
    user_council_models = request.council_models
    user_chairman_model = request.chairman_model
    web_search_enabled = request.web_search_enabled
    speed_mode = request.speed_mode

    # ── Speed Mode parameters ──────────────────────────────────────
    # When speed_mode is enabled, reduce timeouts and cap token output
    # to make every stage respond faster.
    SPEED_TIMEOUT = 60.0 if speed_mode else 120.0
    SPEED_S1_MAX_TOKENS = 2048 if speed_mode else None
    SPEED_S2_MAX_TOKENS = 1536 if speed_mode else None
    SPEED_S3_MAX_TOKENS = 4096 if speed_mode else None

    # ── Stage 2 prompt builder ─────────────────────────────────────
    def _build_stage2_prompt(is_speed: bool) -> str:
        """Return the Stage 2 ranking prompt.

        In speed mode the claim-classification Part 2 is skipped entirely,
        cutting the generated output roughly in half and significantly
        reducing model latency.
        """
        if is_speed:
            return """You are a pharmaceutical domain expert evaluating different responses to the following question:
{context_note}
Question: {question}

Here are the responses from different models (anonymized):

{responses_text}

═══════════════════════════════════════════════════════════
PART 1 — RUBRIC EVALUATION
═══════════════════════════════════════════════════════════
For EACH response, provide a score from 0 to 10 on each criterion below.
After each score, give a brief justification (1 sentence).

Criteria:
  • Relevancy (0-10): How directly and completely the response addresses the original question
  • Faithfulness (0-10): Factual accuracy, absence of hallucinations
  • Context Recall (0-10): Coverage of key concepts and nuances
  • Output Quality (0-10): Clarity, structure, depth
  • Consensus (0-10): Would other domain experts broadly agree?

Format EXACTLY as follows for each response:

RUBRIC Response X:
  Relevancy: <score>/10 — <justification>
  Faithfulness: <score>/10 — <justification>
  Context Recall: <score>/10 — <justification>
  Output Quality: <score>/10 — <justification>
  Consensus: <score>/10 — <justification>

═══════════════════════════════════════════════════════════
FINAL RANKING
═══════════════════════════════════════════════════════════
Based on your rubric evaluation above, provide your final ranking
from best to worst.

IMPORTANT: Your final ranking MUST be formatted EXACTLY as follows:
- Start with the line "FINAL RANKING:" (all caps, with colon)
- Then list the responses from best to worst as a numbered list
- Each line should be: number, period, space, then ONLY the response label (e.g., "1. Response A")

Now provide your complete evaluation:"""
        # ── Full prompt (default) ──────────────────────────────────
        return """You are a pharmaceutical domain expert evaluating different responses to the following question:
{context_note}
Question: {question}

Here are the responses from different models (anonymized):

{responses_text}

═══════════════════════════════════════════════════════════
PART 1 — RUBRIC EVALUATION (Verbalized Sampling)
═══════════════════════════════════════════════════════════
For EACH response, provide a score from 0 to 10 on each criterion below.
After each score, give a brief justification (1-2 sentences).

Criteria:
  • Relevancy (0-10): How directly and completely the response addresses the original question
  • Faithfulness (0-10): Factual accuracy, absence of hallucinations, grounded in evidence
  • Context Recall (0-10): Coverage of key concepts, dimensions, and nuances raised across all responses
  • Output Quality (0-10): Clarity, structure, depth, readability, and overall coherence
  • Consensus (0-10): Would other domain experts broadly agree with the claims made?

Format EXACTLY as follows for each response:

RUBRIC Response X:
  Relevancy: <score>/10 — <justification>
  Faithfulness: <score>/10 — <justification>
  Context Recall: <score>/10 — <justification>
  Output Quality: <score>/10 — <justification>
  Consensus: <score>/10 — <justification>

═══════════════════════════════════════════════════════════
PART 2 — CLAIM CLASSIFICATION (Pharma Safety)
═══════════════════════════════════════════════════════════
For EACH response, classify its major claims in pharmaceutical context:
  TP (True Positive)  = Correct, verifiable claim relevant to the question
  FP (False Positive) = Incorrect, misleading, or hallucinated claim
  FN (False Negative) = Important information the response FAILED to mention

Format EXACTLY as follows for each response:

CLAIMS Response X:
  TP: <count> — <brief summary of correct claims>
  FP: <count> — <brief summary of incorrect/hallucinated claims, or "None detected">
  FN: <count> — <brief summary of important omissions, or "None detected">

═══════════════════════════════════════════════════════════
PART 3 — FINAL RANKING
═══════════════════════════════════════════════════════════
Based on your rubric evaluation and claim analysis above, provide
your final ranking from best to worst.

IMPORTANT: Your final ranking MUST be formatted EXACTLY as follows:
- Start with the line "FINAL RANKING:" (all caps, with colon)
- Then list the responses from best to worst as a numbered list
- Each line should be: number, period, space, then ONLY the response label (e.g., "1. Response A")

Now provide your complete evaluation:"""

    # Generate a unique session ID for kill switch tracking
    session_id = f"{conversation_id}:{uuid.uuid4().hex[:8]}"

    # ── Keep-Alive wrapper ──────────────────────────────────────────
    # Corporate proxies (Zscaler, Netskope) often kill idle SSE
    # connections after ~60-120s.  We send SSE comment pings every
    # 10s to keep the TCP socket alive during long model calls.

    async def with_keepalive(inner_gen, interval=10):
        """
        Wrap an async generator so that a ': keepalive\\n\\n' SSE
        comment is emitted whenever `interval` seconds elapse between
        real data frames.

        Uses asyncio.create_task so that the inner generator's
        __anext__() is never cancelled — asyncio.wait_for would
        cancel it on timeout, corrupting the generator state.

        Context Propagation:
        The inner generator (event_generator) sets ContextVars like
        _current_memory_user during its first __anext__() invocation.
        Since ensure_future() creates a new Task for each __anext__()
        with a fresh context copy, those ContextVars would be lost.
        We capture the context from the first Task (which ran the
        generator's setup code), then reuse it for all subsequent
        Tasks so that ContextVars persist across yields.
        """
        inner_iter = inner_gen.__aiter__()
        pending_next = None          # the Task for __anext__()
        gen_context = None           # captured from first Task for propagation

        while True:
            if pending_next is None:
                coro = inner_iter.__anext__()
                if gen_context is not None:
                    # Reuse the generator's context so ContextVars
                    # (e.g. _current_memory_user) persist across yields
                    loop = asyncio.get_running_loop()
                    pending_next = loop.create_task(coro, context=gen_context)
                else:
                    pending_next = asyncio.ensure_future(coro)

            done, _ = await asyncio.wait({pending_next}, timeout=interval)

            if done:
                # The inner generator yielded a value (or raised)
                task = pending_next
                pending_next = None

                # After the first iteration, capture the Task's context
                # which includes ContextVars set by the generator (e.g.
                # set_memory_user).  All future Tasks reuse this context.
                if gen_context is None:
                    gen_context = task.get_context()

                try:
                    yield task.result()          # may raise StopAsyncIteration
                except StopAsyncIteration:
                    break
            else:
                # Timeout — inner gen is still working; send keepalive
                yield ": keepalive\n\n"

    async def event_generator():
        nonlocal augmented_content
        # Register with kill switch
        kill_event = kill_switch.register_session(session_id)
        cost_tracker = SessionCostTracker()
        # Scope all memory operations to this user
        set_memory_user(user_id)
        try:
            # Emit session ID so the frontend can target the kill switch
            yield f"data: {json.dumps({'type': 'session_start', 'data': {'session_id': session_id}})}\n\n"

            # ── Pipeline Timer ──────────────────────────────────
            timer = PipelineTimer()
            timer.start("total")

            # Check global halt before starting
            if kill_switch.is_halted:
                yield f"data: {json.dumps({'type': 'error', 'message': 'System is in emergency halt mode. Please try again later.', 'code': 'GLOBAL_HALT'})}\n\n"
                return

            # ── Prompt Suitability Guard ─────────────────────────────
            # Block unsuitable prompts before any stage is triggered.
            # If conversation was already blocked, reject all follow-ups.
            if conversation.get("blocked"):
                yield f"data: {json.dumps({'type': 'prompt_rejected', 'data': {'category': 'CONVERSATION_BLOCKED', 'message': 'This conversation has been closed due to a policy violation. Please start a **new conversation** with a question related to pharmaceutical sciences, clinical research, or life-science topics.'}})}\n\n"
                return

            # ── OPT-1: Prompt Guard ∥ Pre-Stage 1 Memory Recall ──────
            # Fire both in parallel — memory recall is read-only so safe
            # to discard if the guard rejects the prompt.
            # Use augmented_content (which includes extracted file text)
            # so the guard's keyword bank and LLM check can see the
            # actual pharma/science content from attached documents.
            guard_input = augmented_content or request.content or ""
            timer.start("prompt_guard")
            guard_task = asyncio.create_task(evaluate_prompt(guard_input, has_attachments=bool(request.attachments)))
            memory_task = asyncio.create_task(pre_stage1_agent(augmented_content, conversation_id, user_id=user_id))

            guard_verdict = await guard_task
            timer.stop("prompt_guard")
            if not guard_verdict.allowed:
                memory_task.cancel()  # discard memory work
                # Mark conversation as blocked so no follow-ups are accepted
                conversation["blocked"] = True
                conversation["blocked_reason"] = guard_verdict.category
                storage.save_conversation(user_id, conversation)
                storage_content_early = request.content or ""
                storage.add_user_message(user_id, conversation_id, storage_content_early)

                # ── Still classify & title even for rejected prompts ──
                # This ensures the sidebar shows a meaningful title and
                # domain tag instead of "New Conversation" with no label.
                try:
                    from .memory import UserProfileMemory
                    guard_context_tags = UserProfileMemory.classify_query(guard_input)
                    storage.update_conversation_context(user_id, conversation_id, guard_context_tags)
                    yield f"data: {json.dumps({'type': 'context_classified', 'data': guard_context_tags})}\n\n"
                except Exception:
                    pass  # non-fatal
                try:
                    guard_title = await generate_conversation_title(storage_content_early)
                    if guard_title:
                        storage.update_conversation_title(user_id, conversation_id, guard_title)
                        yield f"data: {json.dumps({'type': 'title_complete', 'data': {'title': guard_title}})}\n\n"
                except Exception:
                    pass  # non-fatal

                yield f"data: {json.dumps({'type': 'prompt_rejected', 'data': {'category': guard_verdict.category, 'message': guard_verdict.message}})}\n\n"
                return

            # Guard passed — await memory recall result
            timer.start("memory_recall")
            memory_gate = await memory_task
            timer.stop("memory_recall")
            raw_memory_context = memory_gate.get("memory_context", "")
            duplicate_episode = None
            if memory_gate.get("memory_context"):
                augmented_content = memory_gate["augmented_query"]
            if memory_gate.get("duplicate_detected") and memory_gate.get("duplicate_episode"):
                duplicate_episode = {
                    **memory_gate["duplicate_episode"],
                    "duplicate_similarity": memory_gate.get("duplicate_similarity", 0),
                }
            memory_recall_data = {k: v for k, v in memory_gate.items() if k != 'memory_context' and k != 'augmented_query'}
            yield f"data: {json.dumps({'type': 'memory_recall', 'data': memory_recall_data})}\n\n"

            # Build user message content for storage (include attachment info)
            storage_content = request.content
            if request.attachments:
                attachment_list = [f"📎 {att.name}" for att in request.attachments]
                if storage_content:
                    storage_content = f"{storage_content}\n\n---\nAttachments:\n" + "\n".join(attachment_list)
                else:
                    storage_content = "Attachments:\n" + "\n".join(attachment_list)
            
            # Get existing conversation history for follow-up context
            conversation_history = conversation.get("messages", [])
            
            # Add user message
            storage.add_user_message(user_id, conversation_id, storage_content)

            # ── Targeted Follow-Up Detection ─────────────────────────
            # When the user clicks a FOCUS ON chip (Stage X or model name),
            # the frontend prepends "Regarding Stage X: " or "Regarding
            # <model>'s response: " to the user's question.  In these cases
            # we skip the full 3-stage pipeline and instead send the
            # referenced content directly to the chairman for a focused
            # answer.  Agent team still fires for analysis.
            _targeted_followup = _detect_targeted_followup(
                augmented_content, conversation_history
            )
            if _targeted_followup:
                logger.info(
                    f"[Targeted Follow-Up] type={_targeted_followup['type']}, "
                    f"target={_targeted_followup.get('target_label', 'N/A')}"
                )
                # Yield the targeted follow-up path (much faster)
                async for event in _run_targeted_followup(
                    targeted=_targeted_followup,
                    user_query=augmented_content,
                    conversation_history=conversation_history,
                    conversation_id=conversation_id,
                    user_id=user_id,
                    user_chairman_model=user_chairman_model,
                    user_council_models=user_council_models,
                    web_search_enabled=web_search_enabled,
                    speed_mode=speed_mode,
                    session_id=session_id,
                    SPEED_TIMEOUT=SPEED_TIMEOUT,
                    SPEED_S3_MAX_TOKENS=SPEED_S3_MAX_TOKENS,
                    cost_tracker=cost_tracker,
                ):
                    yield event
                return  # exit event_generator — targeted path is complete

            # Start title generation in parallel (don't await yet)
            title_task = None
            if is_first_message:
                title_task = asyncio.create_task(generate_conversation_title(request.content))

            # Stage 1: Collect responses with incremental progress
            # Instead of waiting for all models to finish before emitting
            # any data, we fire individual tasks and yield progress events
            # as each model completes — eliminating the multi-minute dead zone.
            timer.start("stage1")
            yield f"data: {json.dumps({'type': 'stage1_start'})}\n\n"

            from .config import COUNCIL_MODELS as _COUNCIL_MODELS

            s1_models = user_council_models or _COUNCIL_MODELS
            s1_context = build_conversation_context(conversation_history)
            s1_query = f"{s1_context}Current question (follow-up): {augmented_content}" if s1_context else augmented_content
            s1_messages = [{"role": "user", "content": s1_query}]

            # Kill switch gate
            if session_id and kill_switch.is_session_killed(session_id):
                raise KillSwitchError(f"Session {session_id} killed before Stage 1")

            # Fire all model queries as individual tasks (max_retries=1
            # since the self-healing fallback mechanism already substitutes
            # a different model on failure — no need for 2 retries + 4.5s
            # of backoff delay per failing model).
            pending_tasks = {}
            _s1_model_starts = {}  # track per-model start times
            for model in s1_models:
                task = asyncio.create_task(
                    query_model(model, s1_messages,
                                timeout=SPEED_TIMEOUT,
                                web_search_enabled=web_search_enabled,
                                session_id=session_id,
                                max_retries=1,
                                max_tokens=SPEED_S1_MAX_TOKENS)
                )
                pending_tasks[task] = model
                _s1_model_starts[model] = time.perf_counter()

            stage1_results = []
            s1_failed_models = []
            s1_used_models = set(s1_models)
            s1_total = len(s1_models)

            while pending_tasks:
                done, _ = await asyncio.wait(
                    pending_tasks.keys(), return_when=asyncio.FIRST_COMPLETED
                )
                for task in done:
                    model = pending_tasks.pop(task)
                    try:
                        response = task.result()
                        _s1_model_ms = round((time.perf_counter() - _s1_model_starts.get(model, 0)) * 1000, 1)
                        timer.record_model("stage1", model, _s1_model_ms)
                        if response is not None:
                            result_item = {
                                "model": model,
                                "response": response.get('content', ''),
                                "usage": response.get('usage'),
                            }
                            stage1_results.append(result_item)
                            yield f"data: {json.dumps({'type': 'stage1_model_complete', 'data': result_item, 'progress': {'completed': len(stage1_results), 'failed': len(s1_failed_models), 'total': s1_total}})}\n\n"
                        else:
                            s1_failed_models.append(model)
                    except Exception as e:
                        logger.error(f"[Stage1] {model} raised: {e}")
                        circuit_breaker.record_failure(model, str(e))
                        s1_failed_models.append(model)

            # Self-healing: attempt fallback models for any that failed
            if s1_failed_models and len(stage1_results) < len(s1_models):
                logger.info(f"[Stage1] {len(s1_failed_models)} model(s) failed, attempting fallbacks...")
                for failed_model in s1_failed_models:
                    if session_id and kill_switch.is_session_killed(session_id):
                        raise KillSwitchError(f"Session {session_id} killed during Stage 1 fallback")
                    fallback = resolve_fallback(failed_model, s1_used_models)
                    if fallback:
                        s1_used_models.add(fallback)
                        fb_response = await query_model(
                            fallback, s1_messages,
                            web_search_enabled=web_search_enabled,
                            session_id=session_id,
                        )
                        if fb_response is not None:
                            result_item = {
                                "model": f"{fallback} (fallback for {failed_model})",
                                "response": fb_response.get('content', ''),
                                "usage": fb_response.get('usage'),
                            }
                            stage1_results.append(result_item)
                            health_monitor.log_healing_action("stage1_fallback_success", {
                                "failed_model": failed_model, "fallback_model": fallback,
                            })
                            yield f"data: {json.dumps({'type': 'stage1_model_complete', 'data': result_item, 'progress': {'completed': len(stage1_results), 'failed': len(s1_failed_models), 'total': s1_total}})}\n\n"
                        else:
                            health_monitor.log_healing_action("stage1_fallback_failed", {
                                "failed_model": failed_model, "fallback_model": fallback,
                            })

            # Quorum check
            if not check_quorum(stage1_results, "Stage 1", MIN_STAGE1_QUORUM):
                health_monitor.log_healing_action("stage1_quorum_failure", {
                    "successful": len(stage1_results), "required": MIN_STAGE1_QUORUM,
                })
                raise QuorumError(
                    f"Stage 1 quorum not met: got {len(stage1_results)}, need {MIN_STAGE1_QUORUM}"
                )

            # Record Stage 1 token usage
            for r in stage1_results:
                cost_tracker.record("stage1", r["model"], r.get("usage"))
            timer.stop("stage1")
            yield f"data: {json.dumps({'type': 'stage1_complete', 'data': stage1_results})}\n\n"

            # ── Checkpoint: Stage 1 complete ──────────────────────────
            try:
                storage.save_pipeline_checkpoint(user_id, conversation_id, {
                    "completed_stage": "stage1",
                    "stage1_results": stage1_results,
                    "augmented_content": augmented_content,
                    "web_search_enabled": web_search_enabled,
                    "user_council_models": user_council_models,
                    "user_chairman_model": user_chairman_model,
                    "conversation_history": conversation_history,
                    "raw_memory_context": raw_memory_context if 'raw_memory_context' in dir() else None,
                })
            except Exception as e:
                logger.debug(f"[Checkpoint] Non-fatal save error after Stage 1: {e}")

            # ── Early Title Save + Context Classification ─────────
            # Save title and context tags RIGHT AFTER Stage 1 so the
            # sidebar updates immediately and memory/skills can index
            # the conversation by domain even if the pipeline fails.
            context_tags = None
            if title_task:
                try:
                    timer.start("title_generation")
                    title = await title_task
                    timer.stop("title_generation")
                    title_task = None  # consumed
                    storage.update_conversation_title(user_id, conversation_id, title)
                    yield f"data: {json.dumps({'type': 'title_complete', 'data': {'title': title}})}\n\n"
                except Exception as e:
                    timer.stop("title_generation")
                    logger.warning(f"[EarlyTitle] Non-fatal: {e}")

            # Classify the query for context tagging (domain, type, complexity)
            try:
                from .memory import UserProfileMemory
                timer.start("context_classify")
                context_tags = UserProfileMemory.classify_query(augmented_content)
                timer.stop("context_classify")
                storage.update_conversation_context(user_id, conversation_id, context_tags)
                yield f"data: {json.dumps({'type': 'context_classified', 'data': context_tags})}\n\n"
            except Exception as e:
                timer.stop("context_classify")
                logger.warning(f"[ContextTags] Non-fatal: {e}")

            # Kill switch check between stages
            if kill_switch.is_session_killed(session_id):
                yield f"data: {json.dumps({'type': 'killed', 'message': 'Council session aborted by user.'})}\n\n"
                return

            # ── OPT-3: Incremental Stage 2 — stream each ranking as it arrives ──
            # Also fires evidence retrieval in parallel with Stage 2.
            timer.start("stage2")
            yield f"data: {json.dumps({'type': 'stage2_start'})}\n\n"
            timer.start("evidence_retrieval")
            evidence_task = asyncio.create_task(run_evidence_skills(augmented_content, web_search_enabled=web_search_enabled))

            # Build Stage 2 prompt (same logic as stage2_collect_rankings but inline)
            _s2_models = user_council_models or COUNCIL_MODELS
            _s2_labels = [chr(65 + i) for i in range(len(stage1_results))]
            label_to_model = {
                f"Response {label}": result['model']
                for label, result in zip(_s2_labels, stage1_results)
            }
            _s2_context = build_conversation_context(conversation_history)
            _s2_context_note = f"\nNote: This is a follow-up question in an ongoing conversation.\n{_s2_context}\n" if _s2_context else ""
            _s2_responses_text = "\n\n".join([
                f"Response {label}:\n{result['response']}"
                for label, result in zip(_s2_labels, stage1_results)
            ])

            # ── Position Debiasing (arXiv:2405.19323) ──────────────
            # Shuffle presentation order PER REVIEWER to mitigate
            # first-position bias.  Labels stay the same — only the
            # display order changes.
            _s2_per_model_responses: dict[str, str] = {}
            for _s2m in _s2_models:
                _shuf = list(range(len(stage1_results)))
                random.shuffle(_shuf)
                _s2_per_model_responses[_s2m] = "\n\n".join([
                    f"Response {_s2_labels[i]}:\n{stage1_results[i]['response']}"
                    for i in _shuf
                ])

            _s2_ranking_prompt_template = _build_stage2_prompt(speed_mode)

            # Build per-model debiased messages
            _s2_per_model_msgs: dict[str, list] = {}
            for _s2m in _s2_models:
                _s2_per_model_msgs[_s2m] = [{"role": "user", "content": _s2_ranking_prompt_template.format(
                    context_note=_s2_context_note,
                    question=augmented_content,
                    responses_text=_s2_per_model_responses.get(_s2m, _s2_responses_text),
                )}]

            _s2_tasks = {
                asyncio.create_task(
                    query_model(m, _s2_per_model_msgs[m],
                                timeout=SPEED_TIMEOUT,
                                web_search_enabled=web_search_enabled,
                                session_id=session_id,
                                max_tokens=SPEED_S2_MAX_TOKENS)
                ): m
                for m in _s2_models
            }
            _s2_model_starts = {m: time.perf_counter() for m in _s2_models}
            stage2_results = []
            _s2_pending = set(_s2_tasks.keys())
            _s2_total = len(_s2_models)

            while _s2_pending:
                done, _s2_pending = await asyncio.wait(_s2_pending, return_when=asyncio.FIRST_COMPLETED)
                for task in done:
                    model_name = _s2_tasks[task]
                    try:
                        resp = task.result()
                        _s2_model_ms = round((time.perf_counter() - _s2_model_starts.get(model_name, 0)) * 1000, 1)
                        timer.record_model("stage2", model_name, _s2_model_ms)
                        if resp is not None:
                            full_text = resp.get('content', '')
                            parsed = parse_ranking_from_text(full_text)
                            rubric = parse_rubric_scores(full_text)
                            claims = parse_claim_counts(full_text)
                            result_item = {
                                "model": model_name,
                                "ranking": full_text,
                                "parsed_ranking": parsed,
                                "rubric_scores": rubric,
                                "claim_counts": claims,
                                "usage": resp.get('usage'),
                            }
                            stage2_results.append(result_item)
                            cost_tracker.record("stage2", model_name, resp.get('usage'))
                            yield f"data: {json.dumps({'type': 'stage2_model_response', 'data': result_item, 'progress': {'completed': len(stage2_results), 'total': _s2_total}})}\n\n"
                    except Exception as e:
                        logger.error(f"[Stage2-Incremental] {model_name} failed: {e}")

            if len(stage2_results) < MIN_STAGE2_QUORUM:
                raise QuorumError(f"Stage 2 quorum not met: got {len(stage2_results)}, need {MIN_STAGE2_QUORUM}")

            aggregate_rankings = calculate_aggregate_rankings(stage2_results, label_to_model)
            # Compute grounding scores
            timer.start("grounding_compute")
            grounding_scores = compute_response_grounding_scores(
                stage2_results, label_to_model, aggregate_rankings
            )
            timer.stop("grounding_compute")
            # Await evidence retrieval (should be done by now — ran during Stage 2)
            evidence_bundle = await evidence_task
            timer.stop("evidence_retrieval")
            yield f"data: {json.dumps({'type': 'evidence_complete', 'data': evidence_bundle})}\n\n"
            timer.stop("stage2")
            yield f"data: {json.dumps({'type': 'stage2_complete', 'data': stage2_results, 'metadata': {'label_to_model': label_to_model, 'aggregate_rankings': aggregate_rankings, 'grounding_scores': grounding_scores}})}\n\n"

            # ── Checkpoint: Stage 2 complete ──────────────────────────
            try:
                storage.save_pipeline_checkpoint(user_id, conversation_id, {
                    "completed_stage": "stage2",
                    "stage1_results": stage1_results,
                    "stage2_results": stage2_results,
                    "label_to_model": label_to_model,
                    "aggregate_rankings": aggregate_rankings,
                    "grounding_scores": grounding_scores,
                    "augmented_content": augmented_content,
                    "web_search_enabled": web_search_enabled,
                    "user_council_models": user_council_models,
                    "user_chairman_model": user_chairman_model,
                    "conversation_history": conversation_history,
                    "raw_memory_context": raw_memory_context if 'raw_memory_context' in dir() else None,
                })
            except Exception as e:
                logger.debug(f"[Checkpoint] Non-fatal save error after Stage 2: {e}")

            # ── Stage 2.5: Relevancy Gate ────────────────────────────
            relevancy_gate = compute_relevancy_gate(stage2_results)
            gated_labels = [lbl for lbl, g in relevancy_gate.items() if g.get("gated_out")]
            yield f"data: {json.dumps({'type': 'relevancy_gate', 'data': {'gate': relevancy_gate, 'gated_labels': gated_labels}})}\n\n"

            # Kill switch check between stages
            if kill_switch.is_session_killed(session_id):
                yield f"data: {json.dumps({'type': 'killed', 'message': 'Council session aborted by user.'})}\n\n"
                return

            # ── OPT-4: Post-Stage 2 ∥ Stage 3 Start ─────────────────
            # Stage 3 does NOT depend on post_stage2_agent output.
            # Fire Stage 3 + CA validation + post_stage2 ALL in parallel
            # so the chairman LLM call starts immediately.
            yield f"data: {json.dumps({'type': 'stage3_start'})}\n\n"
            evidence_text = format_citations_for_prompt(evidence_bundle)

            # Fire CA validation pass in parallel with Stage 3 (lightweight probes)
            # Skip in speed mode to reduce latency
            ca_validation_task = None
            if not speed_mode:
                ca_validation_task = asyncio.create_task(
                    stage2_ca_validation_pass(
                        augmented_content,
                        stage1_results,
                        label_to_model,
                        user_council_models,
                        web_search_enabled,
                        session_id=session_id,
                    )
                )

            # Fire Stage 3 via token-by-token streaming for faster TTFB
            # Build Stage 3 prompt (pure computation, no API call)
            stage3_messages, chairman_to_use = build_stage3_prompt(
                user_query=augmented_content,
                stage1_results=stage1_results,
                stage2_results=stage2_results,
                chairman_model=user_chairman_model,
                conversation_history=conversation_history,
                evidence_context=evidence_text,
                relevancy_gate=relevancy_gate,
                memory_context=raw_memory_context,
                duplicate_episode=duplicate_episode,
            )

            # Fire post-Stage 2 agent in parallel (doesn't block Stage 3)
            post_stage2_task = asyncio.create_task(
                post_stage2_agent(
                    augmented_content, grounding_scores, aggregate_rankings,
                    user_id=user_id,
                )
            )

            # Await post-Stage 2 (fast memory lookup, emits SSE while Stage 3 runs)
            stage2_gate = await post_stage2_task
            yield f"data: {json.dumps({'type': 'memory_gate', 'data': stage2_gate})}\n\n"

            # Kill switch check (cancel running tasks if killed)
            if kill_switch.is_session_killed(session_id):
                if ca_validation_task:
                    ca_validation_task.cancel()
                yield f"data: {json.dumps({'type': 'killed', 'message': 'Council session aborted by user.'})}\n\n"
                return

            # ── Stream Stage 3 tokens ────────────────────────────────
            streamed_text = ""
            stage3_usage = None
            stream_succeeded = False
            try:
                async for chunk in query_model_stream(
                    chairman_to_use,
                    stage3_messages,
                    timeout=SPEED_TIMEOUT or 150.0,
                    web_search_enabled=web_search_enabled,
                    session_id=session_id,
                    max_tokens=SPEED_S3_MAX_TOKENS,
                ):
                    if isinstance(chunk, dict):
                        # Final sentinel with usage info
                        stage3_usage = chunk.get("usage")
                    else:
                        streamed_text += chunk
                        yield f"data: {json.dumps({'type': 'stage3_chunk', 'data': {'text': chunk}})}\n\n"
                if streamed_text:
                    stream_succeeded = True
                    timer.stop("stage3_streaming")
                    logger.info(f"[Stage3] Streaming complete — {len(streamed_text)} chars from {chairman_to_use}")
            except Exception as e:
                timer.stop("stage3_streaming")
                logger.warning(f"[Stage3] Streaming failed ({e}), falling back to non-streaming")

            # If streaming produced no text, fall back to non-streaming with speculative racing
            if stream_succeeded:
                stage3_result = {
                    "model": chairman_to_use,
                    "response": streamed_text,
                    "usage": stage3_usage,
                }
            else:
                timer.start("stage3_fallback")
                stage3_result = await stage3_synthesize_final(
                    augmented_content, stage1_results, stage2_results,
                    user_chairman_model, conversation_history, web_search_enabled,
                    session_id=session_id,
                    evidence_context=evidence_text,
                    relevancy_gate=relevancy_gate,
                    memory_context=raw_memory_context,
                    duplicate_episode=duplicate_episode,
                    max_tokens=SPEED_S3_MAX_TOKENS,
                    timeout=SPEED_TIMEOUT,
                )
                timer.stop("stage3_fallback")
                logger.info(f"[Stage3] Fallback completed — {stage3_result.get('model', 'unknown')}")

            # Record Stage 3 token usage
            cost_tracker.record("stage3", stage3_result.get("model", "unknown"), stage3_result.get("usage"))

            # ── OPT-7: Fire DT immediately ∥ CA processing ─────────
            # DT does NOT depend on CA results — fire its LLM call now
            # so it overlaps with CA post-processing (saves 5-15s).
            yield f"data: {json.dumps({'type': 'doubting_thomas_start'})}\n\n"
            timer.start("doubting_thomas")
            dt_task = asyncio.create_task(
                doubting_thomas_review(
                    user_query=augmented_content,
                    draft_response=stage3_result.get("response", ""),
                    stage1_results=stage1_results,
                    relevancy_gate=relevancy_gate,
                    chairman_model=user_chairman_model,
                    web_search_enabled=web_search_enabled,
                    session_id=session_id,
                )
            )

            # While DT runs, process CA validation (should already be complete from Stage 3)
            # Skipped in speed mode (ca_validation_task is None)
            if ca_validation_task is not None:
                try:
                    timer.start("ca_validation")
                    ca_validation_results = await ca_validation_task
                    timer.stop("ca_validation")
                    if ca_validation_results:
                        # Record CA validation token usage
                        for model_name, val_data in ca_validation_results.items():
                            if val_data.get("usage"):
                                cost_tracker.record("ca_validation", model_name, val_data["usage"])
                        # Enhance grounding scores with multi-round CA data
                        grounding_scores = enhance_ca_with_validation(
                            grounding_scores, ca_validation_results, label_to_model
                        )
                        yield f"data: {json.dumps({'type': 'ca_validation_complete', 'data': {'models_probed': len(ca_validation_results), 'grounding_scores': grounding_scores}})}\n\n"
                        logger.info(f"[CA Validation] Enhanced CA for {len(ca_validation_results)} models")

                        # Persist CA snapshots for cross-session tracking
                        try:
                            set_memory_user(user_id)  # re-set after yield boundary
                            memory_mgr = get_memory_manager()
                            for resp in grounding_scores.get("per_response", []):
                                ca = resp.get("context_awareness")
                                if ca and ca.get("score") is not None:
                                    memory_mgr.store_ca_snapshot(
                                        conversation_id=conversation_id,
                                        model=resp["model"],
                                        ca_data=ca,
                                    )
                        except Exception as e:
                            logger.debug(f"[CA Tracking] Non-fatal persistence error: {e}")
                except Exception as e:
                    timer.stop("ca_validation")
                    logger.warning(f"[CA Validation] Non-fatal error: {e}")
                    # Continue without enhanced CA — original grounding_scores remain

            # ── Await Doubting Thomas (detect-and-fix self-reflection) ──
            # (arXiv:2602.03837 §Adversarial Reviewer; arXiv:2602.13949 §Reflection)
            # ALWAYS runs — even in speed mode — to preserve quality assurance
            dt_result = None
            try:
                dt_result = await dt_task
                if dt_result.get("fix_applied"):
                    stage3_result["response"] = dt_result["revised_response"]
                    # Record DT token usage
                    cost_tracker.record("doubting_thomas", user_chairman_model or "chairman", dt_result.get("usage"))
                    logger.info(
                        f"[Doubting Thomas] Revised synthesis applied — "
                        f"{dt_result['defect_count']} defect(s) fixed"
                    )
                timer.stop("doubting_thomas")
                yield f"data: {json.dumps({'type': 'doubting_thomas_complete', 'data': {'defect_count': dt_result.get('defect_count', 0), 'needs_fix': dt_result.get('needs_fix', False), 'fix_applied': dt_result.get('fix_applied', False), 'criteria': dt_result.get('criteria', []), 'fix_instructions': dt_result.get('fix_instructions', []), 'critique': dt_result.get('critique') or None}})}\n\n"
            except Exception as e:
                timer.stop("doubting_thomas")
                logger.warning(f"[Doubting Thomas] Non-fatal error: {e}")
                yield f"data: {json.dumps({'type': 'doubting_thomas_complete', 'data': {'defect_count': 0, 'needs_fix': False, 'fix_applied': False, 'criteria': [], 'fix_instructions': [], 'error': str(e)}})}\n\n"

            # Extract infographic data from the chairman's response
            infographic_data = None
            raw_response = stage3_result.get("response", "")
            infographic_data = extract_infographic(raw_response)
            # Strip the raw infographic JSON block from the displayed response
            stage3_result["response"] = strip_infographic_block(raw_response)

            # ── Citation Supervisor: enrich references with clickable links ──
            timer.start("citation_enrich")
            stage3_result["response"] = enrich_stage3_citations(
                stage3_result["response"]
            )
            timer.stop("citation_enrich")

            # ── Citation Validator: verify URLs are reachable, fix broken ones ──
            # Skipped in speed mode (URL validation is I/O-heavy)
            if not speed_mode:
                try:
                    timer.start("citation_validate")
                    stage3_result["response"] = await validate_and_fix_citations(
                        stage3_result["response"]
                    )
                    timer.stop("citation_validate")
                except Exception as e:
                    timer.stop("citation_validate")
                    logger.warning(f"[Citation Validator] Non-fatal error: {e}")

            yield f"data: {json.dumps({'type': 'stage3_complete', 'data': stage3_result})}\n\n"

            # Emit infographic data as a separate event
            if infographic_data:
                yield f"data: {json.dumps({'type': 'infographic_complete', 'data': infographic_data})}\n\n"

            # Title was saved early (after Stage 1). If it wasn't
            # consumed yet (e.g. Stage 1 had zero results path), save now.
            if title_task:
                try:
                    title = await title_task
                    storage.update_conversation_title(user_id, conversation_id, title)
                    yield f"data: {json.dumps({'type': 'title_complete', 'data': {'title': title}})}\n\n"
                except Exception:
                    pass  # non-critical fallback

            # Save complete assistant message (including metadata for reload)
            try:
                storage.add_assistant_message(
                    user_id,
                    conversation_id,
                    stage1_results,
                    stage2_results,
                    stage3_result,
                    metadata={
                        "label_to_model": label_to_model,
                        "aggregate_rankings": aggregate_rankings,
                        "grounding_scores": grounding_scores,
                        "evidence": evidence_bundle,
                        "infographic": infographic_data,
                        "memory_recall": memory_recall_data,
                        "memory_gate": stage2_gate,
                        "relevancy_gate": relevancy_gate,
                        "doubting_thomas": {
                            "defect_count": dt_result.get("defect_count", 0) if dt_result else 0,
                            "needs_fix": dt_result.get("needs_fix", False) if dt_result else False,
                            "fix_applied": dt_result.get("fix_applied", False) if dt_result else False,
                        },
                        "context_tags": context_tags,
                    },
                )
            except Exception as e:
                logger.error(f"[Storage] Failed to save assistant message: {e}")
                yield f"data: {json.dumps({'type': 'storage_warning', 'data': {'message': 'Results displayed but could not be saved to history.'}})}\n\n"

            # Emit cost summary before completion
            timer.stop("total")
            timing_summary = timer.summary()
            cost_summary = cost_tracker.compute_summary()
            cost_summary["timing"] = timing_summary
            # Attach Redis cache stats if available
            try:
                from .memory_store import get_redis_stats, _get_redis_client
                if _get_redis_client() is not None:
                    cost_summary["redis_cache"] = get_redis_stats()
            except Exception:
                pass
            yield f"data: {json.dumps({'type': 'cost_summary', 'data': cost_summary})}\n\n"

            # ── OPT-5: Agent Team ∥ Post-Stage 3 Learning ─────────
            # Both are independent of each other — fire in parallel.
            overall_grounding = grounding_scores.get("overall_score", 0) / 100.0

            timer.start("agent_team")
            timer.start("learning")
            agent_team_task = asyncio.create_task(
                run_agent_team(
                    user_query=augmented_content,
                    stage1_results=stage1_results,
                    stage2_results=stage2_results,
                    stage3_result=stage3_result,
                    aggregate_rankings=aggregate_rankings,
                    grounding_scores=grounding_scores,
                    evidence_bundle=evidence_bundle,
                    cost_summary=cost_summary,
                    web_search_enabled=web_search_enabled,
                )
            )
            learning_task = asyncio.create_task(
                post_stage3_agent(
                    conversation_id=conversation_id,
                    user_query=augmented_content,
                    stage1_results=stage1_results,
                    aggregate_rankings=aggregate_rankings,
                    stage3_result=stage3_result,
                    grounding_score=overall_grounding,
                    cost_summary=cost_summary,
                    tags=[context_tags["domain"], context_tags.get("question_type", "general")] if context_tags else None,
                    user_id=user_id,
                )
            )

            # Await both (whichever finishes first proceeds; both are fast)
            agent_team_result = None
            learning_gate = None
            try:
                agent_team_result, learning_gate = await asyncio.gather(
                    agent_team_task, learning_task, return_exceptions=True
                )
            except Exception as e:
                logger.warning(f"[Parallel post-pipeline] gather error: {e}")
            timer.stop("agent_team")
            timer.stop("learning")

            # Emit agent team result
            if agent_team_result and not isinstance(agent_team_result, Exception):
                try:
                    yield f"data: {json.dumps({'type': 'agent_team_complete', 'data': agent_team_result})}\n\n"
                    storage.update_last_message_metadata(
                        user_id, conversation_id, {"agent_team": agent_team_result}
                    )
                except Exception:
                    logger.debug("Failed to persist agent_team metadata")
            elif isinstance(agent_team_result, Exception):
                logger.warning(f"[AgentTeam] Non-fatal error: {agent_team_result}")

            # Emit learning decision
            if learning_gate and not isinstance(learning_gate, Exception):
                yield f"data: {json.dumps({'type': 'memory_learning', 'data': learning_gate})}\n\n"
                # Persist memory_learning + cost_summary in conversation metadata
                try:
                    storage.update_last_message_metadata(
                        user_id, conversation_id,
                        {"memory_learning": learning_gate, "cost_summary": cost_summary},
                    )
                except Exception:
                    logger.debug("Failed to persist memory_learning metadata")
            elif isinstance(learning_gate, Exception):
                logger.warning(f"[PostStage3] Non-fatal error: {learning_gate}")
                # Still persist cost_summary even if learning failed
                try:
                    storage.update_last_message_metadata(
                        user_id, conversation_id, {"cost_summary": cost_summary},
                    )
                except Exception:
                    logger.debug("Failed to persist cost_summary metadata")

            # ── OPT-6: User Behaviour Learning + ECA Adaptation ───
            # Non-blocking: record user profile and run adaptation loop.
            try:
                set_memory_user(user_id)  # re-set after yield boundary
                upm = get_user_profile_memory()
                eca = get_eca()

                # Classify the query and record the interaction
                classification = upm.classify_query(augmented_content)
                upm.record_interaction(
                    user_id=user_id,
                    query=augmented_content,
                    grounding_score=overall_grounding,
                    relevancy_violations=gated_labels,
                    gated_labels=gated_labels,
                    classification=classification,
                )

                # Get updated user profile for ECA
                user_profile = upm.get_user_profile(user_id)

                # Run full ECA adaptation (Memory × Skills pairing)
                eca_result = eca.run_full_adaptation(
                    user_id=user_id,
                    user_profile=user_profile,
                    evidence_bundle=evidence_bundle,
                    grounding_scores=grounding_scores,
                    grounding_score_overall=overall_grounding,
                )

                # Emit user behaviour + ECA event
                yield f"data: {json.dumps({'type': 'user_behaviour_update', 'data': {'profile': user_profile, 'eca': eca_result, 'classification': classification}})}\n\n"

                # Persist ECA state in conversation metadata
                storage.update_last_message_metadata(
                    user_id, conversation_id,
                    {"user_profile": user_profile, "eca_adaptation": eca_result, "relevancy_gate": relevancy_gate},
                )
                logger.info(
                    f"[UserBehaviour+ECA] Complete for {user_id}: "
                    f"domain={classification['domain']}, ema_reward={eca_result.get('ema_reward', 0):.3f}"
                )
            except Exception as e:
                logger.warning(f"[UserBehaviour+ECA] Non-fatal error: {e}")

            # Send completion event
            yield f"data: {json.dumps({'type': 'complete'})}\n\n"

            # ── Clear pipeline checkpoint on successful completion ──
            try:
                storage.clear_pipeline_checkpoint(user_id, conversation_id)
            except Exception:
                pass  # non-critical

        except KillSwitchError as e:
            # User-triggered abort — graceful termination
            yield f"data: {json.dumps({'type': 'killed', 'message': str(e)})}\n\n"
        except QuorumError as e:
            # Self-healing exhausted — not enough models responded
            yield f"data: {json.dumps({'type': 'error', 'message': f'Self-healing exhausted: {e}', 'code': 'QUORUM_FAILURE'})}\n\n"
        except Exception as e:
            # Send error event
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"
        finally:
            # Always unregister the session from kill switch
            kill_switch.unregister_session(session_id)

    return StreamingResponse(
        with_keepalive(event_generator(), interval=10),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache, no-store, must-revalidate",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
            "Pragma": "no-cache",
            "Expires": "0",
        }
    )


# ────────────────────────────────────────────────────────────────────────
# Resume endpoint — self-healing SSE reconnect after network drop
# ────────────────────────────────────────────────────────────────────────

class ResumeRequest(BaseModel):
    """Request body for the resume endpoint."""
    council_models: Optional[List[str]] = None
    chairman_model: Optional[str] = None
    web_search_enabled: bool = False
    speed_mode: bool = False


@app.post("/api/conversations/{conversation_id}/message/resume")
async def resume_message_stream(
    conversation_id: str,
    request: ResumeRequest,
    user_id: str = Depends(get_authenticated_user_id),
):
    """Resume an interrupted council pipeline from its last checkpoint.

    Reads the progressive checkpoint saved after Stage 1 / Stage 2 and
    re-runs only the remaining stages, returning SSE events for the
    uncompleted portion of the pipeline.
    """
    conversation = storage.get_conversation(user_id, conversation_id)
    if conversation is None:
        raise HTTPException(status_code=404, detail="Conversation not found")

    checkpoint = storage.load_pipeline_checkpoint(user_id, conversation_id)
    if not checkpoint:
        raise HTTPException(
            status_code=409,
            detail="No pipeline checkpoint found — please resend the message.",
        )

    completed_stage = checkpoint.get("completed_stage")
    if completed_stage not in ("stage1", "stage2"):
        raise HTTPException(
            status_code=409,
            detail=f"Unexpected checkpoint stage: {completed_stage}",
        )

    session_id = f"{conversation_id}:{uuid.uuid4().hex[:8]}"

    # Reuse keepalive from the main endpoint's closure — define inline
    async def with_keepalive_resume(inner_gen, interval=10):
        inner_iter = inner_gen.__aiter__()
        pending_next = None
        gen_context = None
        while True:
            if pending_next is None:
                coro = inner_iter.__anext__()
                if gen_context is not None:
                    loop = asyncio.get_running_loop()
                    pending_next = loop.create_task(coro, context=gen_context)
                else:
                    pending_next = asyncio.ensure_future(coro)
            done, _ = await asyncio.wait({pending_next}, timeout=interval)
            if done:
                task = pending_next
                pending_next = None
                if gen_context is None:
                    gen_context = task.get_context()
                try:
                    yield task.result()
                except StopAsyncIteration:
                    break
            else:
                yield ": keepalive\n\n"

    async def resume_generator():
        kill_event = kill_switch.register_session(session_id)
        cost_tracker = SessionCostTracker()
        set_memory_user(user_id)

        try:
            yield f"data: {json.dumps({'type': 'session_start', 'data': {'session_id': session_id}})}\n\n"

            # ── Restore checkpoint data ─────────────────────────────
            stage1_results = checkpoint["stage1_results"]
            augmented_content = checkpoint["augmented_content"]
            web_search_enabled = checkpoint.get("web_search_enabled", request.web_search_enabled)
            user_council_models = checkpoint.get("user_council_models", request.council_models)
            user_chairman_model = checkpoint.get("user_chairman_model", request.chairman_model)
            conversation_history = checkpoint.get("conversation_history", [])
            raw_memory_context = checkpoint.get("raw_memory_context")

            # Classify the resumed query for context tagging
            resume_context_tags = None
            try:
                from .memory import UserProfileMemory
                resume_context_tags = UserProfileMemory.classify_query(augmented_content)
                storage.update_conversation_context(user_id, conversation_id, resume_context_tags)
                yield f"data: {json.dumps({'type': 'context_classified', 'data': resume_context_tags})}\n\n"
            except Exception as e:
                logger.warning(f"[Resume-ContextTags] Non-fatal: {e}")

            # Re-emit Stage 1 data so the frontend hydrates its state
            yield f"data: {json.dumps({'type': 'stage1_complete', 'data': stage1_results})}\n\n"

            if completed_stage == "stage1":
                # ── Need to run Stage 2 + Stage 3 ──────────────────
                yield f"data: {json.dumps({'type': 'stage2_start'})}\n\n"
                evidence_task = asyncio.create_task(
                    run_evidence_skills(augmented_content, web_search_enabled=web_search_enabled)
                )

                _s2_models = user_council_models or COUNCIL_MODELS
                _s2_labels = [chr(65 + i) for i in range(len(stage1_results))]
                label_to_model = {
                    f"Response {label}": result['model']
                    for label, result in zip(_s2_labels, stage1_results)
                }
                _s2_context = build_conversation_context(conversation_history)
                _s2_context_note = (
                    f"\nNote: This is a follow-up question in an ongoing conversation.\n{_s2_context}\n"
                    if _s2_context else ""
                )
                _s2_responses_text = "\n\n".join([
                    f"Response {label}:\n{result['response']}"
                    for label, result in zip(_s2_labels, stage1_results)
                ])

                # Position debiasing
                _s2_per_model_responses: dict[str, str] = {}
                for _s2m in _s2_models:
                    _shuf = list(range(len(stage1_results)))
                    random.shuffle(_shuf)
                    _s2_per_model_responses[_s2m] = "\n\n".join([
                        f"Response {_s2_labels[i]}:\n{stage1_results[i]['response']}"
                        for i in _shuf
                    ])

                _s2_ranking_prompt_template = """You are a pharmaceutical domain expert evaluating different responses to the following question:
{context_note}
Question: {question}

Here are the responses from different models (anonymized):

{responses_text}

\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550
PART 1 \u2014 RUBRIC EVALUATION (Verbalized Sampling)
\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550
For EACH response, provide a score from 0 to 10 on each criterion below.
After each score, give a brief justification (1-2 sentences).

Criteria:
  \u2022 Relevancy (0-10): How directly and completely the response addresses the original question
  \u2022 Faithfulness (0-10): Factual accuracy, absence of hallucinations, grounded in evidence
  \u2022 Context Recall (0-10): Coverage of key concepts, dimensions, and nuances raised across all responses
  \u2022 Output Quality (0-10): Clarity, structure, depth, readability, and overall coherence
  \u2022 Consensus (0-10): Would other domain experts broadly agree with the claims made?

Format EXACTLY as follows for each response:

RUBRIC Response X:
  Relevancy: <score>/10 \u2014 <justification>
  Faithfulness: <score>/10 \u2014 <justification>
  Context Recall: <score>/10 \u2014 <justification>
  Output Quality: <score>/10 \u2014 <justification>
  Consensus: <score>/10 \u2014 <justification>

\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550
PART 2 \u2014 CLAIM CLASSIFICATION (Pharma Safety)
\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550
For EACH response, classify its major claims in pharmaceutical context:
  TP (True Positive)  = Correct, verifiable claim relevant to the question
  FP (False Positive) = Incorrect, misleading, or hallucinated claim
  FN (False Negative) = Important information the response FAILED to mention

Format EXACTLY as follows for each response:

CLAIMS Response X:
  TP: <count> \u2014 <brief summary of correct claims>
  FP: <count> \u2014 <brief summary of incorrect/hallucinated claims, or "None detected">
  FN: <count> \u2014 <brief summary of important omissions, or "None detected">

\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550
PART 3 \u2014 FINAL RANKING
\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550
Based on your rubric evaluation and claim analysis above, provide
your final ranking from best to worst.

IMPORTANT: Your final ranking MUST be formatted EXACTLY as follows:
- Start with the line "FINAL RANKING:" (all caps, with colon)
- Then list the responses from best to worst as a numbered list
- Each line should be: number, period, space, then ONLY the response label (e.g., "1. Response A")

Now provide your complete evaluation:"""

                _s2_per_model_msgs: dict[str, list] = {}
                for _s2m in _s2_models:
                    _s2_per_model_msgs[_s2m] = [{"role": "user", "content": _s2_ranking_prompt_template.format(
                        context_note=_s2_context_note,
                        question=augmented_content,
                        responses_text=_s2_per_model_responses.get(_s2m, _s2_responses_text),
                    )}]

                _s2_tasks = {
                    asyncio.create_task(
                        query_model(m, _s2_per_model_msgs[m], web_search_enabled=web_search_enabled, session_id=session_id)
                    ): m
                    for m in _s2_models
                }
                stage2_results = []
                _s2_pending = set(_s2_tasks.keys())
                _s2_total = len(_s2_models)

                while _s2_pending:
                    done, _s2_pending = await asyncio.wait(_s2_pending, return_when=asyncio.FIRST_COMPLETED)
                    for task in done:
                        model_name = _s2_tasks[task]
                        try:
                            resp = task.result()
                            if resp is not None:
                                full_text = resp.get('content', '')
                                parsed = parse_ranking_from_text(full_text)
                                rubric = parse_rubric_scores(full_text)
                                claims = parse_claim_counts(full_text)
                                result_item = {
                                    "model": model_name,
                                    "ranking": full_text,
                                    "parsed_ranking": parsed,
                                    "rubric_scores": rubric,
                                    "claim_counts": claims,
                                    "usage": resp.get('usage'),
                                }
                                stage2_results.append(result_item)
                                cost_tracker.record("stage2", model_name, resp.get('usage'))
                                yield f"data: {json.dumps({'type': 'stage2_model_response', 'data': result_item, 'progress': {'completed': len(stage2_results), 'total': _s2_total}})}\n\n"
                        except Exception as e:
                            logger.error(f"[Resume-Stage2] {model_name} failed: {e}")

                if len(stage2_results) < MIN_STAGE2_QUORUM:
                    raise QuorumError(f"Stage 2 quorum not met: got {len(stage2_results)}, need {MIN_STAGE2_QUORUM}")

                aggregate_rankings = calculate_aggregate_rankings(stage2_results, label_to_model)
                grounding_scores = compute_response_grounding_scores(stage2_results, label_to_model, aggregate_rankings)
                evidence_bundle = await evidence_task
                yield f"data: {json.dumps({'type': 'evidence_complete', 'data': evidence_bundle})}\n\n"
                yield f"data: {json.dumps({'type': 'stage2_complete', 'data': stage2_results, 'metadata': {'label_to_model': label_to_model, 'aggregate_rankings': aggregate_rankings, 'grounding_scores': grounding_scores}})}\n\n"

                # Save Stage 2 checkpoint
                try:
                    storage.save_pipeline_checkpoint(user_id, conversation_id, {
                        "completed_stage": "stage2",
                        "stage1_results": stage1_results,
                        "stage2_results": stage2_results,
                        "label_to_model": label_to_model,
                        "aggregate_rankings": aggregate_rankings,
                        "grounding_scores": grounding_scores,
                        "augmented_content": augmented_content,
                        "web_search_enabled": web_search_enabled,
                        "user_council_models": user_council_models,
                        "user_chairman_model": user_chairman_model,
                        "conversation_history": conversation_history,
                        "raw_memory_context": raw_memory_context,
                    })
                except Exception:
                    pass

            elif completed_stage == "stage2":
                # Stage 2 already done — restore its results
                stage2_results = checkpoint["stage2_results"]
                label_to_model = checkpoint["label_to_model"]
                aggregate_rankings = checkpoint["aggregate_rankings"]
                grounding_scores = checkpoint["grounding_scores"]

                # Re-emit Stage 2 data for frontend hydration
                yield f"data: {json.dumps({'type': 'stage2_complete', 'data': stage2_results, 'metadata': {'label_to_model': label_to_model, 'aggregate_rankings': aggregate_rankings, 'grounding_scores': grounding_scores}})}\n\n"

                # Run evidence retrieval (fast, needed for Stage 3)
                evidence_bundle = await run_evidence_skills(augmented_content, web_search_enabled=web_search_enabled)
                yield f"data: {json.dumps({'type': 'evidence_complete', 'data': evidence_bundle})}\n\n"

            # ── Relevancy Gate ──────────────────────────────────────
            relevancy_gate = compute_relevancy_gate(stage2_results)
            gated_labels = [lbl for lbl, g in relevancy_gate.items() if g.get("gated_out")]
            yield f"data: {json.dumps({'type': 'relevancy_gate', 'data': {'gate': relevancy_gate, 'gated_labels': gated_labels}})}\n\n"

            # ── Stage 3 (streaming) ────────────────────────────────
            yield f"data: {json.dumps({'type': 'stage3_start'})}\n\n"
            evidence_text = format_citations_for_prompt(evidence_bundle)

            # Build prompt, stream tokens, fallback to non-streaming
            stage3_messages_tf, chairman_to_use_tf = build_stage3_prompt(
                user_query=augmented_content,
                stage1_results=stage1_results,
                stage2_results=stage2_results,
                chairman_model=user_chairman_model,
                conversation_history=conversation_history,
                evidence_context=evidence_text,
                relevancy_gate=relevancy_gate,
                memory_context=raw_memory_context,
            )

            streamed_text_tf = ""
            stage3_usage_tf = None
            stream_ok_tf = False
            try:
                async for chunk in query_model_stream(
                    chairman_to_use_tf,
                    stage3_messages_tf,
                    timeout=SPEED_TIMEOUT or 150.0,
                    web_search_enabled=web_search_enabled,
                    session_id=session_id,
                    max_tokens=SPEED_S3_MAX_TOKENS,
                ):
                    if isinstance(chunk, dict):
                        stage3_usage_tf = chunk.get("usage")
                    else:
                        streamed_text_tf += chunk
                        yield f"data: {json.dumps({'type': 'stage3_chunk', 'data': {'text': chunk}})}\n\n"
                if streamed_text_tf:
                    stream_ok_tf = True
            except Exception as e:
                logger.warning(f"[Stage3-TF] Streaming failed ({e}), falling back")

            if stream_ok_tf:
                stage3_result = {
                    "model": chairman_to_use_tf,
                    "response": streamed_text_tf,
                    "usage": stage3_usage_tf,
                }
            else:
                stage3_result = await stage3_synthesize_final(
                    augmented_content, stage1_results, stage2_results,
                    user_chairman_model, conversation_history, web_search_enabled,
                    session_id=session_id,
                    evidence_context=evidence_text,
                    relevancy_gate=relevancy_gate,
                    memory_context=raw_memory_context,
                    max_tokens=SPEED_S3_MAX_TOKENS,
                    timeout=SPEED_TIMEOUT,
                )
            cost_tracker.record("stage3", stage3_result.get("model", "unknown"), stage3_result.get("usage"))

            # ── Doubting Thomas ─────────────────────────────────────
            # Skipped in speed mode
            dt_result = None
            if not speed_mode:
                try:
                    yield f"data: {json.dumps({'type': 'doubting_thomas_start'})}\n\n"
                    dt_result = await doubting_thomas_review(
                        user_query=augmented_content,
                        draft_response=stage3_result.get("response", ""),
                        stage1_results=stage1_results,
                        relevancy_gate=relevancy_gate,
                        chairman_model=user_chairman_model,
                        web_search_enabled=web_search_enabled,
                        session_id=session_id,
                    )
                    if dt_result.get("fix_applied"):
                        stage3_result["response"] = dt_result["revised_response"]
                        cost_tracker.record("doubting_thomas", user_chairman_model or "chairman", dt_result.get("usage"))
                    yield f"data: {json.dumps({'type': 'doubting_thomas_complete', 'data': {'defect_count': dt_result.get('defect_count', 0), 'needs_fix': dt_result.get('needs_fix', False), 'fix_applied': dt_result.get('fix_applied', False), 'criteria': dt_result.get('criteria', []), 'fix_instructions': dt_result.get('fix_instructions', []), 'critique': dt_result.get('critique') or None}})}\n\n"
                except Exception as e:
                    logger.warning(f"[Resume-DT] Non-fatal error: {e}")
                    yield f"data: {json.dumps({'type': 'doubting_thomas_complete', 'data': {'defect_count': 0, 'needs_fix': False, 'fix_applied': False, 'criteria': [], 'fix_instructions': [], 'error': str(e)}})}\n\n"
            else:
                yield f"data: {json.dumps({'type': 'doubting_thomas_complete', 'data': {'defect_count': 0, 'needs_fix': False, 'fix_applied': False, 'skipped': True}})}\n\n"

            # ── Citation enrichment ──────────────────────────────────
            raw_response = stage3_result.get("response", "")
            infographic_data = extract_infographic(raw_response)
            stage3_result["response"] = strip_infographic_block(raw_response)
            stage3_result["response"] = enrich_stage3_citations(stage3_result["response"])
            if not speed_mode:
                try:
                    stage3_result["response"] = await validate_and_fix_citations(stage3_result["response"])
                except Exception:
                    pass

            yield f"data: {json.dumps({'type': 'stage3_complete', 'data': stage3_result})}\n\n"

            if infographic_data:
                yield f"data: {json.dumps({'type': 'infographic_complete', 'data': infographic_data})}\n\n"

            # ── Save final assistant message ────────────────────────
            storage.add_assistant_message(
                user_id, conversation_id,
                stage1_results, stage2_results, stage3_result,
                metadata={
                    "label_to_model": label_to_model,
                    "aggregate_rankings": aggregate_rankings,
                    "grounding_scores": grounding_scores,
                    "evidence": evidence_bundle,
                    "infographic": infographic_data,
                    "relevancy_gate": relevancy_gate,
                    "context_tags": resume_context_tags,
                    "doubting_thomas": {
                        "defect_count": dt_result.get("defect_count", 0) if dt_result else 0,
                        "needs_fix": dt_result.get("needs_fix", False) if dt_result else False,
                        "fix_applied": dt_result.get("fix_applied", False) if dt_result else False,
                    },
                },
            )

            # ── Cost summary ────────────────────────────────────────
            cost_summary = cost_tracker.compute_summary()
            yield f"data: {json.dumps({'type': 'cost_summary', 'data': cost_summary})}\n\n"

            # ── Agent Team ──────────────────────────────────────────
            try:
                agent_team_result = await run_agent_team(
                    user_query=augmented_content,
                    stage1_results=stage1_results,
                    stage2_results=stage2_results,
                    stage3_result=stage3_result,
                    aggregate_rankings=aggregate_rankings,
                    grounding_scores=grounding_scores,
                    evidence_bundle=evidence_bundle,
                    cost_summary=cost_summary,
                    web_search_enabled=web_search_enabled,
                )
                if agent_team_result:
                    yield f"data: {json.dumps({'type': 'agent_team_complete', 'data': agent_team_result})}\n\n"
                    storage.update_last_message_metadata(user_id, conversation_id, {"agent_team": agent_team_result})
            except Exception as e:
                logger.warning(f"[Resume-AgentTeam] Non-fatal error: {e}")

            # ── Clear checkpoint & complete ──────────────────────────
            try:
                storage.clear_pipeline_checkpoint(user_id, conversation_id)
            except Exception:
                pass

            yield f"data: {json.dumps({'type': 'complete'})}\n\n"

        except KillSwitchError as e:
            yield f"data: {json.dumps({'type': 'killed', 'message': str(e)})}\n\n"
        except QuorumError as e:
            yield f"data: {json.dumps({'type': 'error', 'message': f'Self-healing exhausted: {e}', 'code': 'QUORUM_FAILURE'})}\n\n"
        except Exception as e:
            logger.error(f"[Resume] Error: {e}", exc_info=True)
            yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"
        finally:
            kill_switch.unregister_session(session_id)

    return StreamingResponse(
        with_keepalive_resume(resume_generator(), interval=10),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache, no-store, must-revalidate",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
            "Pragma": "no-cache",
            "Expires": "0",
        },
    )


if __name__ == "__main__":
    import uvicorn

    ssl_certfile = os.getenv("SSL_CERTFILE")
    ssl_keyfile = os.getenv("SSL_KEYFILE")

    uvicorn_kwargs = dict(host="0.0.0.0", port=8001)
    if ssl_certfile and ssl_keyfile:
        uvicorn_kwargs["ssl_certfile"] = ssl_certfile
        uvicorn_kwargs["ssl_keyfile"] = ssl_keyfile
        uvicorn_kwargs["ssl_version"] = 2  # TLS 1.3 (ssl.TLS_VERSION_TLSv1_3)
        print("\n🔒 TLS ENABLED — serving over HTTPS")
    else:
        print("\n⚠️  No SSL_CERTFILE/SSL_KEYFILE — serving over plain HTTP (dev mode)")

    uvicorn.run(app, **uvicorn_kwargs)
