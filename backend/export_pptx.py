"""
PPTX export v3 — every slide gets contextual AI images.

Uses the multi-provider image_gen module for parallel image generation.
Section slides: full-width 16:9 hero images.
Content slides: portrait image strip (left) + text card (right).
"""

import io
import re
from typing import List, Dict, Any, Optional

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

from .image_gen import (
    generate_images_parallel_sync,
    build_slide_prompt,
    build_section_prompt,
)

# Brand colours
BAYER_GREEN = RGBColor(0x10, 0x85, 0x7F)
BAYER_DARK = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)
GRAY_TEXT = RGBColor(0x55, 0x65, 0x81)

# Stage themes
STAGE_THEME = {
    1: {'bg': RGBColor(0xE7, 0xF3, 0xFF), 'accent': RGBColor(0x15, 0x65, 0xC0), 'title': 'Stage 1 — Model Responses'},
    2: {'bg': RGBColor(0xFF, 0xF4, 0xE6), 'accent': RGBColor(0xE6, 0x51, 0x00), 'title': 'Stage 2 — Peer Rankings'},
    3: {'bg': RGBColor(0xE8, 0xF5, 0xE9), 'accent': RGBColor(0x2E, 0x7D, 0x32), 'title': "Stage 3 — Chairman Synthesis"},
}

# Layout constants (16:9 slide)
SLIDE_WIDTH = Cm(33.87)
SLIDE_HEIGHT = Cm(19.05)
MARGIN = Cm(1.2)
MAX_CHARS = 1200

# Content slide split layout
IMG_STRIP_WIDTH = Cm(9.5)
CONTENT_LEFT = Cm(11.5)
CONTENT_WIDTH = Cm(21.0)


def _strip_md(text: str) -> str:
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    text = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'', text)
    text = re.sub(r'^#{1,6}\s*', '', text, flags=re.MULTILINE)
    return text.strip()


def _chunk(text: str, limit: int = MAX_CHARS) -> List[str]:
    if len(text) <= limit:
        return [text]
    out, buf, n = [], [], 0
    for line in text.split('\n'):
        ln = len(line) + 1
        if n + ln > limit and buf:
            out.append('\n'.join(buf))
            buf, n = [line], ln
        else:
            buf.append(line)
            n += ln
    if buf:
        out.append('\n'.join(buf))
    return out or [text]


def _add_title_slide(prs: Presentation, conversation: Dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # Light background — NO black/dark backgrounds
    bg = slide.shapes.add_shape(1, 0, 0, sw, sh)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()

    band = slide.shapes.add_shape(1, 0, Cm(7), sw, Cm(0.6))
    band.fill.solid(); band.fill.fore_color.rgb = BAYER_GREEN; band.line.fill.background()

    tb = slide.shapes.add_textbox(Cm(2.5), Cm(3.0), sw - Cm(5.0), Cm(3.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = conversation.get('title', 'LLM Council Report')
    p.font.size = Pt(34); p.font.bold = True
    p.font.color.rgb = BAYER_DARK; p.alignment = PP_ALIGN.CENTER
    sp = tf.add_paragraph(); sp.text = 'Multi-Model Deliberation'
    sp.font.size = Pt(16); sp.font.color.rgb = BAYER_GREEN; sp.alignment = PP_ALIGN.CENTER
    if conversation.get('created_at'):
        dp = tf.add_paragraph()
        dp.text = f"Generated: {conversation['created_at']}"
        dp.font.size = Pt(11)
        dp.font.color.rgb = GRAY_TEXT
        dp.alignment = PP_ALIGN.CENTER

    footer = slide.shapes.add_textbox(Cm(2.0), sh - Cm(1.5), sw - Cm(4.0), Cm(1.0))
    ft = footer.text_frame
    fp = ft.paragraphs[0]; fp.text = 'Powered by: llmcouncil@bayer.com team members'
    fp.font.size = Pt(10)
    fp.font.color.rgb = GRAY_TEXT
    fp.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs: Presentation, stage: int, headline: str,
                       kicker: str = '', hero_image: Optional[bytes] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    theme = STAGE_THEME.get(stage, STAGE_THEME[1])
    sw, sh = prs.slide_width, prs.slide_height

    top = slide.shapes.add_shape(1, 0, 0, sw, Cm(3.5))
    top.fill.solid(); top.fill.fore_color.rgb = theme['accent']; top.line.fill.background()
    body = slide.shapes.add_shape(1, 0, Cm(3.5), sw, sh - Cm(3.5))
    body.fill.solid(); body.fill.fore_color.rgb = theme['bg']; body.line.fill.background()

    tb = slide.shapes.add_textbox(MARGIN, Cm(1.2), sw - 2 * MARGIN, Cm(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = headline
    p.font.size = Pt(26); p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); p.alignment = PP_ALIGN.LEFT
    if kicker:
        kp = tf.add_paragraph(); kp.text = kicker
        kp.font.size = Pt(14)
        kp.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        kp.alignment = PP_ALIGN.LEFT

    if hero_image:
        try:
            stream = io.BytesIO(hero_image)
            img_width = sw - 2 * MARGIN
            slide.shapes.add_picture(stream, MARGIN, Cm(3.9), width=img_width)
        except Exception:
            pass


def _add_content_slide(prs: Presentation, title: str, body: str,
                       accent: RGBColor, stage: int,
                       image: Optional[bytes] = None):
    """Add a content slide with optional portrait image strip on the left."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # Top accent bar
    bar = slide.shapes.add_shape(1, 0, 0, sw, Cm(0.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    content_top = Cm(0.8)
    content_height = sh - Cm(1.4)

    has_image = False
    if image:
        try:
            stream = io.BytesIO(image)
            slide.shapes.add_picture(
                stream, MARGIN, content_top,
                width=IMG_STRIP_WIDTH, height=content_height,
            )
            has_image = True
        except Exception:
            pass

    if has_image:
        card_left = CONTENT_LEFT
        card_width = CONTENT_WIDTH
    else:
        card_left = MARGIN
        card_width = sw - 2 * MARGIN

    # White card background
    rect = slide.shapes.add_shape(1, card_left, content_top, card_width, content_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    rect.line.color.rgb = accent
    rect.line.width = Pt(1.5)

    # Text content
    tb = slide.shapes.add_textbox(
        card_left + Cm(0.4), content_top + Cm(0.35),
        card_width - Cm(0.8), content_height - Cm(0.7),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    h = tf.paragraphs[0]
    h.text = title
    h.font.size = Pt(14)
    h.font.bold = True
    h.font.color.rgb = accent

    for line in body.split('\n'):
        p = tf.add_paragraph()
        clean = line.strip()
        if not clean:
            p.text = ''
            p.space_after = Pt(4)
            continue
        if clean.startswith('- '):
            p.text = '\u2022 ' + clean[2:]
            p.font.color.rgb = accent
        elif re.match(r'^\d+[\.\)]\s', clean):
            p.text = clean
            p.font.color.rgb = BAYER_DARK
        else:
            p.text = clean
            p.font.color.rgb = BAYER_DARK
        p.font.size = Pt(11)
        p.space_after = Pt(2)


def _collect_image_prompts(conversation: Dict[str, Any],
                           user_question: str) -> Dict[str, str]:
    """Collect all image prompts for the entire deck.

    Returns {key: prompt_text}. Keys like 'hero_1', 's1_<model>_c<idx>', etc.
    Each content chunk gets its own unique prompt with a different visual
    perspective to eliminate duplicate imagery across slides.
    """
    prompts: Dict[str, str] = {}

    # Section hero images
    for stage in (1, 2, 3):
        prompts[f'hero_{stage}'] = build_section_prompt(stage, user_question)

    # Per-chunk content images — every content slide gets a unique image
    for msg in conversation.get('messages', []):
        if msg.get('role') == 'user':
            continue

        for resp in (msg.get('stage1') or []):
            model = resp.get('model', 'model')
            raw = _strip_md(resp.get('response', ''))
            chunks = _chunk(raw)
            for ci, chunk_text in enumerate(chunks):
                key = f's1_{model}_c{ci}'
                prompts[key] = build_slide_prompt(
                    model, chunk_text, 1, user_question,
                    chunk_index=ci, total_chunks=len(chunks))

        for ranking in (msg.get('stage2') or []):
            model = ranking.get('model', 'reviewer')
            raw = _strip_md(ranking.get('response', ''))
            chunks = _chunk(raw)
            for ci, chunk_text in enumerate(chunks):
                key = f's2_{model}_c{ci}'
                prompts[key] = build_slide_prompt(
                    model, chunk_text, 2, user_question,
                    chunk_index=ci, total_chunks=len(chunks))

        stage3 = msg.get('stage3')
        if stage3:
            model = stage3.get('model', 'chairman')
            raw = _strip_md(stage3.get('response', ''))
            chunks = _chunk(raw)
            for ci, chunk_text in enumerate(chunks):
                key = f's3_{model}_c{ci}'
                prompts[key] = build_slide_prompt(
                    model, chunk_text, 3, user_question,
                    chunk_index=ci, total_chunks=len(chunks))

    return prompts


def generate_pptx(conversation: Dict[str, Any]) -> bytes:
    """Generate a branded PPTX with contextual images on every slide."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _add_title_slide(prs, conversation)

    # Determine user question
    user_question = ""
    for msg in conversation.get('messages', []):
        if msg.get('role') == 'user' and msg.get('content'):
            user_question = _strip_md(msg.get('content', ''))
            break

    # ── Collect ALL image prompts and generate in parallel ────────
    all_prompts = _collect_image_prompts(conversation, user_question)

    hero_prompts = {k: v for k, v in all_prompts.items() if k.startswith('hero_')}
    content_prompts = {k: v for k, v in all_prompts.items() if not k.startswith('hero_')}

    # Hero images: 16:9 (wide), Content images: 3:4 (portrait sidebar)
    hero_images = generate_images_parallel_sync(
        hero_prompts, aspect_ratio="16:9", max_concurrent=3)
    content_images = generate_images_parallel_sync(
        content_prompts, aspect_ratio="3:4", max_concurrent=4)

    all_images: Dict[str, Optional[bytes]] = {**hero_images, **content_images}

    # ── Build slides ─────────────────────────────────────────────
    for msg in conversation.get('messages', []):
        if msg.get('role') == 'user':
            _add_section_slide(prs, 1, 'User Question', 'What was asked',
                               hero_image=all_images.get('hero_1'))
            _add_content_slide(prs, 'Question',
                               _strip_md(msg.get('content', '')),
                               BAYER_DARK, 1)
            continue

        # Stage 1
        stage1 = msg.get('stage1') or []
        if stage1:
            _add_section_slide(prs, 1, STAGE_THEME[1]['title'],
                               'Key model answers',
                               hero_image=all_images.get('hero_1'))
            for resp in stage1:
                model = resp.get('model', 'Model')
                chunks = _chunk(_strip_md(resp.get('response', '')))
                for idx, chunk in enumerate(chunks, start=1):
                    chunk_img = all_images.get(f's1_{model}_c{idx - 1}')
                    _add_content_slide(
                        prs, f"{model} ({idx}/{len(chunks)})",
                        chunk, STAGE_THEME[1]['accent'], 1,
                        image=chunk_img)

        # Stage 2
        stage2 = msg.get('stage2') or []
        if stage2:
            _add_section_slide(prs, 2, STAGE_THEME[2]['title'],
                               'Peer evaluations',
                               hero_image=all_images.get('hero_2'))
            for ranking in stage2:
                model = ranking.get('model', 'Reviewer')
                chunks = _chunk(_strip_md(ranking.get('response', '')))
                for idx, chunk in enumerate(chunks, start=1):
                    chunk_img = all_images.get(f's2_{model}_c{idx - 1}')
                    _add_content_slide(
                        prs, f"{model} ({idx}/{len(chunks)})",
                        chunk, STAGE_THEME[2]['accent'], 2,
                        image=chunk_img)

        # Stage 3
        stage3 = msg.get('stage3')
        if stage3:
            _add_section_slide(prs, 3, STAGE_THEME[3]['title'],
                               'Final synthesis',
                               hero_image=all_images.get('hero_3'))
            model = stage3.get('model', 'Chairman')
            chunks = _chunk(_strip_md(stage3.get('response', '')))
            for idx, chunk in enumerate(chunks, start=1):
                chunk_img = all_images.get(f's3_{model}_c{idx - 1}')
                _add_content_slide(
                    prs, f"{model} ({idx}/{len(chunks)})",
                    chunk, STAGE_THEME[3]['accent'], 3,
                    image=chunk_img)

    # ── Closing slide (light background) ─────────────────────────
    closing = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    bg = closing.shapes.add_shape(1, 0, 0, sw, sh)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background()
    band = closing.shapes.add_shape(1, 0, Cm(9.0), sw, Cm(0.4))
    band.fill.solid(); band.fill.fore_color.rgb = BAYER_GREEN; band.line.fill.background()
    tb = closing.shapes.add_textbox(Cm(2.0), Cm(5.0), sw - Cm(4.0), Cm(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = 'End of Report'
    p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = BAYER_DARK; p.alignment = PP_ALIGN.CENTER
    sp = tf.add_paragraph(); sp.text = 'LLM Council \u2014 Multi-Model Deliberation'
    sp.font.size = Pt(14); sp.font.color.rgb = BAYER_GREEN; sp.alignment = PP_ALIGN.CENTER
    sp2 = tf.add_paragraph(); sp2.text = 'Powered by: llmcouncil@bayer.com team members'
    sp2.font.size = Pt(11)
    sp2.font.color.rgb = GRAY_TEXT
    sp2.alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
